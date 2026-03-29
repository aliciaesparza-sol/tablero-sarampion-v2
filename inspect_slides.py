from pptx import Presentation
import sys
import io

# Ensure UTF-8 output for terminal
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

file_path = r'c:\Users\aicil\OneDrive\Escritorio\PVU\COEVA\PRESENTACIONES\CAPACITACIÓN SARAMPIÓN CLÍNICA Y VACUNACIÓN 17FEBRERO 2026.pptx'
prs = Presentation(file_path)

for i in range(2, 8):
    if i < len(prs.slides):
        slide = prs.slides[i]
        print(f"--- SLIDE {i} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(f"[{shape.name}] {shape.text}")
        print("-" * 20)
