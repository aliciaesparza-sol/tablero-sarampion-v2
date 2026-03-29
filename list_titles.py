from pptx import Presentation

file_path = r'c:\Users\aicil\OneDrive\Escritorio\PVU\COEVA\PRESENTACIONES\CAPACITACIÓN SARAMPIÓN CLÍNICA Y VACUNACIÓN 17FEBRERO 2026.pptx'
prs = Presentation(file_path)

for i, slide in enumerate(prs.slides):
    title = ""
    if slide.shapes.title:
        title = slide.shapes.title.text
    else:
        # Try to find the first text box that looks like a title
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                title = shape.text.strip().split('\n')[0][:50]
                break
    print(f"Slide {i}: {title}")
