from pptx import Presentation

file_path = r'c:\Users\aicil\OneDrive\Escritorio\PVU\COEVA\PRESENTACIONES\CAPACITACIÓN SARAMPIÓN CLÍNICA Y VACUNACIÓN 17FEBRERO 2026.pptx'
prs = Presentation(file_path)

for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if hasattr(shape, "text") and ("DEFINIC" in shape.text.upper() or "OPERACIONAL" in shape.text.upper()):
            print(f"Match Slide {i}: {shape.text}")
