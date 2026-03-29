from pptx import Presentation

file_path = r'c:\Users\aicil\OneDrive\Escritorio\PVU\COEVA\PRESENTACIONES\CAPACITACIÓN SARAMPIÓN CLÍNICA Y VACUNACIÓN 17FEBRERO 2026.pptx'
prs = Presentation(file_path)

print("Searching for 'DEFINICIÓN' or 'OPERACIONAL'...")
for i, slide in enumerate(prs.slides):
    text = ""
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            text += shape.text.upper() + " "
    
    if "DEFINICIÓN" in text or "DEFINICION" in text or "OPERACIONAL" in text or "CASO" in text:
        print(f"Slide {i}: {text[:150].replace('\n', ' ')}...")
