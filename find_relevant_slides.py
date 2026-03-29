from pptx import Presentation
import sys
import io

# Ensure UTF-8 output for terminal
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

file_path = r'c:\Users\aicil\OneDrive\Escritorio\PVU\COEVA\PRESENTACIONES\CAPACITACIÓN SARAMPIÓN CLÍNICA Y VACUNACIÓN 17FEBRERO 2026.pptx'
prs = Presentation(file_path)

keywords = ["DEFINICIÓN", "OPERACIONAL", "DIFERENCIAL", "SOSP", "PROB", "CONF"]

for i, slide in enumerate(prs.slides):
    text = ""
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            text += shape.text.upper() + " "
    
    found = [k for k in keywords if k in text]
    if found:
        print(f"Slide {i}: Contains {found}")
        print(f"Content snippet: {text[:300]}...")
        print("-" * 20)
