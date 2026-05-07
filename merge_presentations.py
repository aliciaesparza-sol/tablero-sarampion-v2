from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

TEMPLATE_PATH = r"C:\Users\aicil\OneDrive\Escritorio\plantilla presentacion ppt.pptx"
MAP_TERRESTRE = r"C:\Users\aicil\.gemini\antigravity\brain\b90ab7f9-3458-429d-99b9-0806d80e4e60\mapa_rutas_terrestres_mezquital_1778104616998.png"
MAP_AEREO = r"C:\Users\aicil\.gemini\antigravity\brain\b90ab7f9-3458-429d-99b9-0806d80e4e60\mapa_estrategia_aerea_mezquital_1778103476848.png"
OUTPUT_PATH = r"C:\Users\aicil\OneDrive\Escritorio\PVU\SARAMPIÓN\mezquital\PRESENTACIONES\Estrategia_Integrada_Mezquital_COEVA.pptx"

def create_integrated_ppt():
    if not os.path.exists(TEMPLATE_PATH):
        prs = Presentation()
    else:
        prs = Presentation(TEMPLATE_PATH)

    # Helper to add slide
    def add_slide(layout_idx, title_text=None):
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        if title_text and hasattr(slide, "shapes") and slide.shapes.title:
            slide.shapes.title.text = title_text
        return slide

    # Slide 1: Portada
    slide = add_slide(0)
    # Assuming placeholder 0 is title, 1 is subtitle
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            shape.text = "8ª Reunión Extraordinaria del COEVA"
        elif shape.placeholder_format.idx == 1:
            shape.text = "Comité Estatal de Vacunación (COEVA)\nServicios de Salud de Durango\n7 de mayo de 2026"

    # Slide 2: Orden del Día
    slide = add_slide(1, "Orden del Día")
    tf = slide.placeholders[1].text_frame
    items = [
        "Bienvenida y registro.",
        "Propósito de la reunión.",
        "Seguimiento a acuerdos de supervisión federal.",
        "Presentación de estrategia operativa: Terrestre.",
        "Presentación de estrategia operativa: Aérea.",
        "Situación de caravanas estatales/federales.",
        "Acuerdos y compromisos."
    ]
    tf.text = items[0]
    for item in items[1:]:
        p = tf.add_paragraph()
        p.text = item

    # Slide 3: Propósito
    slide = add_slide(1, "Propósito de la Reunión")
    tf = slide.placeholders[1].text_frame
    tf.text = "Continuidad a acuerdos de la supervisión federal (20-24 abril)."
    p = tf.add_paragraph()
    p.text = "Fortalecer la estrategia operativa para Potreros, Cihuacora, San Buenaventura y Curachitos."
    p = tf.add_paragraph()
    p.text = "Articular colaboración CENSIA, INPI, IMSS-Bienestar y SSD."

    # Slide 4: Sección Terrestre
    slide = add_slide(2, "Estrategia Operativa: Vía Terrestre")

    # Slide 5: Mapa Terrestre
    slide = add_slide(5, "Mapa de Rutas Terrestres")
    if os.path.exists(MAP_TERRESTRE):
        prs.slides[-1].shapes.add_picture(MAP_TERRESTRE, Inches(1), Inches(1.5), width=Inches(8))

    # Slide 6: Rutas 1 y 2
    slide = add_slide(1, "Rutas Terrestres 1 y 2")
    tf = slide.placeholders[1].text_frame
    tf.text = "Ruta 1: Durango -> Armadillo (8h). Traslado de remudas para población."
    p = tf.add_paragraph()
    p.text = "Ruta 2: Durango -> Armadillo -> Cihuacora (16h total). Cumple red de frío."

    # Slide 7: Rutas 3 y 4
    slide = add_slide(1, "Rutas Terrestres 3 y 4")
    tf = slide.placeholders[1].text_frame
    tf.text = "Ruta 3: Durango -> Sta. María de Picachos (8h) -> Cihuacora (4h a pie)."
    p = tf.add_paragraph()
    p.text = "Ruta 4: Durango -> Llano Jacalitos (7h) -> Anonas (2h) -> Cihuacora (8h a pie)."

    # Slide 8: Sección Aérea
    slide = add_slide(2, "Estrategia Operativa: Vía Aérea")

    # Slide 9: Mapa Aéreo
    slide = add_slide(5, "Mapa de Operaciones Aéreas")
    if os.path.exists(MAP_AEREO):
        prs.slides[-1].shapes.add_picture(MAP_AEREO, Inches(1), Inches(1.5), width=Inches(8))

    # Slide 10: Plan de Vuelos - Despliegue
    slide = add_slide(1, "Fase 1: Despliegue (18 Mayo)")
    tf = slide.placeholders[1].text_frame
    tf.text = "08:00 AM: Durango -> La Guajolota."
    p = tf.add_paragraph()
    p.text = "Posicionamiento de Equipos PFAM en Potreros y Cihuacora (3 días)."
    p = tf.add_paragraph()
    p.text = "Insumos: Vacunación, red de frío, casa de campaña y medicamentos."

    # Slide 11: Plan de Vuelos - Recuperación
    slide = add_slide(1, "Fase 2: Recuperación (20 Mayo)")
    tf = slide.placeholders[1].text_frame
    tf.text = "13:00 PM: Durango -> Cihuacora."
    p = tf.add_paragraph()
    p.text = "Retorno de PFAM a La Guajolota."
    p = tf.add_paragraph()
    p.text = "ESI continúa en Llanos de Jacalitos y Sombrero Quemado."

    # Slide 12: Caravanas
    slide = add_slide(1, "Estatus de Caravanas")
    tf = slide.placeholders[1].text_frame
    tf.text = "Intervención: 11 al 27 de mayo de 2026."
    p = tf.add_paragraph()
    p.text = "Búsqueda activa de casos de sarampión e inmunización integral."

    # Slide 13: Acuerdos
    slide = add_slide(1, "Acuerdos y Compromisos")
    tf = slide.placeholders[1].text_frame
    tf.text = "Garantizar suficiencia de biológicos y red de frío."
    p = tf.add_paragraph()
    p.text = "Coordinación con INPI para apoyo de remudas en zonas terrestres."
    p = tf.add_paragraph()
    p.text = "Reporte diario de avances al COEVA."

    prs.save(OUTPUT_PATH)
    print(f"Presentación integrada guardada en: {OUTPUT_PATH}")

if __name__ == "__main__":
    create_integrated_ppt()
