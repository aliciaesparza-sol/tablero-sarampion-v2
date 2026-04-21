import os
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE

template_path = r"c:\Users\aicil\OneDrive\Escritorio\plantilla presentacion ppt.pptx"
output_path = r"c:\Users\aicil\OneDrive\Escritorio\Presentacion_Control_Sarampion.pptx"

# Image paths
img_cerco = r"C:\Users\aicil\.gemini\antigravity\brain\a9bcd4dd-ae02-44a1-901d-3dba17160759\cerco_epidemiologico_1776442863511.png"
img_bloqueo = r"C:\Users\aicil\.gemini\antigravity\brain\a9bcd4dd-ae02-44a1-901d-3dba17160759\bloqueo_vacunal_1776442877333.png"
img_barrido = r"C:\Users\aicil\.gemini\antigravity\brain\a9bcd4dd-ae02-44a1-901d-3dba17160759\barrido_documentado_1776442969287.png"
img_monitoreo = r"C:\Users\aicil\.gemini\antigravity\brain\a9bcd4dd-ae02-44a1-901d-3dba17160759\monitoreo_rapido_1776442984217.png"

slides_data = [
    {
        "title": "CERCO EPIDEMIOLÓGICO",
        "bullets": [
            "Verificar cobertura de vacunación (SRP/SR) a la población susceptible; la mención verbal NO es válida, se requiere Cartilla.",
            "Búsqueda activa comunitaria de casos y contactos en torno a los primarios (búsqueda en área y manzanas aledañas).",
            "Inclusión del estudio de la ruta social y crítica: sitios laborales, hospitales, escuelas, medios de transporte público y comunitarios.",
            "Estimular la participación de líderes comunitarios para organizar acciones eficaces por calle, barrio y unidad habitacional."
        ],
        "image": img_cerco
    },
    {
        "title": "BLOQUEO VACUNAL",
        "bullets": [
            "Administración de megadosis de Vitamina A para prevenir complicaciones severas y hospitalizaciones en menores de 5 años.",
            "Búsqueda prioritaria de esquemas nulos o incompletos con vacunas SRP o SR en la población contacto y zonas de riesgo.",
            "Vacunación Urgente: Contactos de 6-11 meses (dosis cero), niños de 1-9 años, y adultos sin antecedente comprobable si pasaron menos de 72 hrs.",
            "Registro inmediato, seguimiento durante estancia domiciliaria y coordinación multisectorial."
        ],
        "image": img_bloqueo
    },
    {
        "title": "BARRIDO DOCUMENTADO",
        "bullets": [
            "Recorridos sistemáticos casa por casa en zonas urbanas (manzanas, edificios, vecindades) y rurales buscando casos y susceptibles.",
            "Manzanas urbanas: iniciar recorrido en una esquina y continuar de manera sistemática hacia la derecha (dirección de manecillas del reloj).",
            "Registro estricto de visitas y de todas las dosis aplicadas llenando el formato completo de censo nominal de vacunación.",
            "Asentamientos irregulares: Identificar un sitio de referencia central o principal para iniciar el recorrido de verificación sistemática."
        ],
        "image": img_barrido
    },
    {
        "title": "MONITOREO RÁPIDO",
        "bullets": [
            "Actualización en tiempo real desde la 'Sala Situacional' evaluando casos probables, confirmados, contactos y tasas de ataque.",
            "Verificación y evaluación sistemática de coberturas en todos los sectores vacunados durante el protocolo del brote epidémico.",
            "La meta institucional es mostrar y garantizar una cobertura mínima del 95% en los territorios y sectores evaluados por las brigadas.",
            "Georreferenciación activa de los casos probables/confirmados, cuadros complicados, defunciones, y de las dosis geolocalizadas aplicadas."
        ],
        "image": img_monitoreo
    }
]

def format_text_frame(tf, bullets):
    tf.clear()  # removes existing content if any
    for b in bullets:
        p = tf.add_paragraph()
        p.text = b
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(20)

def main():
    prs = Presentation(template_path)
    
    # Layout 3 is "Dos objetos"
    layout = prs.slide_layouts[3]
    
    for data in slides_data:
        slide = prs.slides.add_slide(layout)
        
        # Placeholders
        title_ph = slide.placeholders[0]
        text_ph = slide.placeholders[1]
        img_ph = slide.placeholders[2]
        
        # Set Title
        title_ph.text = data["title"]
        for paragraph in title_ph.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(36)
                
        # Set Text Bullets
        format_text_frame(text_ph.text_frame, data["bullets"])
        
        # Insert Image
        # To insert a picture into a placeholder, we retrieve the bounding box
        # of the placeholder shape and add the picture over it, keeping aspect ratio roughly centered.
        # However, a cleaner way is using `img_ph.insert_picture(path)`, provided it's an image or object placeholder.
        try:
            img_ph.insert_picture(data["image"])
        except Exception as e:
            # Fallback if the placeholder doesn't support insert_picture
            print(f"Fallback for image insertion: {e}")
            left = img_ph.left
            top = img_ph.top
            width = img_ph.width
            height = img_ph.height
            slide.shapes.add_picture(data["image"], left, top, width, height)
            
    # Remove any extra empty slides if desired, but we'll just leave the initial template slides
    # OR we can actually remove the ones from the template to make a clean presentation!
    # Let's delete the first slide if it was an empty title from the template, actually
    # since we are creating a presentation *based* on a template, the template might already have slides.
    # To be clean, maybe we just keep what the user had and add our slides at the end.
            
    prs.save(output_path)
    print(f"Saved generated presentation to: {output_path}")

if __name__ == "__main__":
    main()
