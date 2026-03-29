from pptx import Presentation

file_path = r'c:\Users\aicil\OneDrive\Escritorio\PVU\COEVA\PRESENTACIONES\CAPACITACIÓN SARAMPIÓN CLÍNICA Y VACUNACIÓN 17FEBRERO 2026.pptx'
prs = Presentation(file_path)

for i, slide in enumerate(prs.slides):
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            texts.append(shape.text.strip())
    
    full_text = " | ".join(texts)
    print(f"Slide {i}: {full_text[:200]}...")
