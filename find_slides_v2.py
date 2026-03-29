import sys; sys.path.append('.')
from pptx import Presentation

file_path = r'c:\Users\aicil\OneDrive\Escritorio\PVU\COEVA\PRESENTACIONES\CAPACITACIÓN SARAMPIÓN CLÍNICA Y VACUNACIÓN 17FEBRERO 2026.pptx'
prs = Presentation(file_path)

print("Slides Content Summary:")
for i, slide in enumerate(prs.slides):
    text = ""
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            text += shape.text + " "
    title = ""
    if slide.shapes.title:
        title = slide.shapes.title.text
    
    print(f"Slide {i}: Title='{title}' | Content={text[:100].strip()}...")
