import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_improved_presentation(json_file, output_path):
    with open(json_file, 'r', encoding='utf-8') as f:
        slides_info = json.load(f)

    prs = Presentation()
    
    # Define a custom color scheme
    # Primary: Dark Blue (0, 51, 102)
    # Secondary: Teal (0, 128, 128)
    # Background: Light Gray/White
    
    # Slide dimensions (Widescreen 16:9 approx)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1] # Title and Content
    blank_slide_layout = prs.slide_layouts[6]

    for info in slides_info:
        slide_index = info['slide_index']
        text_lines = info['text']
        
        # Filter out empty lines and page numbers
        filtered_text = [t.strip() for t in text_lines if t.strip() and not t.strip().endswith(' / 10')]
        
        if slide_index == 1:
            # Title slide
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Background shape
            bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(240, 248, 255) # Alice Blue
            bg.line.fill.background()
            
            # Header shape
            header = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.5))
            header.fill.solid()
            header.fill.fore_color.rgb = RGBColor(0, 51, 102)
            header.line.fill.background()

            title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(2))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.add_paragraph()
            p.text = filtered_text[0] if len(filtered_text) > 0 else "DIRECTRICES DE VACUNACIÓN"
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 51, 102)
            p.alignment = PP_ALIGN.CENTER
            
            if len(filtered_text) > 1:
                p2 = tf.add_paragraph()
                p2.text = filtered_text[1]
                p2.font.size = Pt(28)
                p2.font.color.rgb = RGBColor(0, 128, 128)
                p2.alignment = PP_ALIGN.CENTER

            if len(filtered_text) > 2:
                sub_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.33), Inches(1.5))
                stf = sub_box.text_frame
                stf.word_wrap = True
                p_sub = stf.add_paragraph()
                p_sub.text = "\n".join(filtered_text[2:])
                p_sub.font.size = Pt(18)
                p_sub.font.color.rgb = RGBColor(100, 100, 100)
                p_sub.alignment = PP_ALIGN.CENTER

        else:
            # Content slide
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Header shape
            header = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.2))
            header.fill.solid()
            header.fill.fore_color.rgb = RGBColor(0, 128, 128)
            header.line.fill.background()
            
            # Title
            title = filtered_text[0] if len(filtered_text) > 0 else "Slide"
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.33), Inches(0.8))
            tf = title_box.text_frame
            p = tf.add_paragraph()
            p.text = title
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            
            # Content
            content_text = filtered_text[1:]
            
            # Try to identify sections (all caps lines)
            y_offset = 1.5
            for line in content_text:
                if line.isupper() and not line.startswith('•'):
                    # Section header
                    box = slide.shapes.add_textbox(Inches(0.5), Inches(y_offset), Inches(12.33), Inches(0.5))
                    p = box.text_frame.add_paragraph()
                    p.text = line
                    p.font.size = Pt(24)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(0, 51, 102)
                    y_offset += 0.6
                else:
                    # Bullet point or normal text
                    box = slide.shapes.add_textbox(Inches(1.0), Inches(y_offset), Inches(11.83), Inches(0.4))
                    p = box.text_frame.add_paragraph()
                    p.text = line
                    p.font.size = Pt(20)
                    p.font.color.rgb = RGBColor(50, 50, 50)
                    y_offset += 0.5
                    
            # Footer
            footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.5))
            fp = footer_box.text_frame.add_paragraph()
            fp.text = "Programa de Vacunación Universal - Lineamientos 2025/2026 (Actualizado con Lineamiento Nacional)"
            fp.font.size = Pt(12)
            fp.font.color.rgb = RGBColor(150, 150, 150)
            fp.alignment = PP_ALIGN.RIGHT

    # Add a concluding slide based on new guidelines
    slide = prs.slides.add_slide(blank_slide_layout)
    header = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(0, 128, 128)
    header.line.fill.background()
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.33), Inches(0.8))
    p = title_box.text_frame.add_paragraph()
    p.text = "Actualización Lineamientos Nacionales 2025-2026"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    content = [
        "ESTRATEGIA INTENSIVA (BROTE)",
        "• Dosis Cero: Aplicación estricta a menores de 6 a 11 meses (no sustituye dosis de 12 meses).",
        "• Población 13-49 años: Vacunación con SR a susceptibles en áreas de alta incidencia.",
        "• Personal de Salud: Refuerzo obligatorio independientemente del esquema previo.",
        "• Bloqueo Vacunal: Vacunación rápida casa por casa ante casos confirmados.",
        "• Validación: Solo cuenta esquema documentado en Cartilla Nacional de Salud."
    ]
    
    y_offset = 1.5
    for line in content:
        if line.isupper() and not line.startswith('•'):
            box = slide.shapes.add_textbox(Inches(0.5), Inches(y_offset), Inches(12.33), Inches(0.5))
            p = box.text_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 51, 102)
            y_offset += 0.6
        else:
            box = slide.shapes.add_textbox(Inches(1.0), Inches(y_offset), Inches(11.83), Inches(0.4))
            p = box.text_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(50, 50, 50)
            y_offset += 0.5
            
    try:
        prs.save(output_path)
        print(f"Successfully saved to {output_path}")
    except Exception as e:
        print(f"Error saving: {e}")

if __name__ == '__main__':
    json_path = r"C:\Users\aicil\.gemini\antigravity\scratch\presentation_content.json"
    output_path = r"c:\Users\aicil\OneDrive\Escritorio\PVU\SARAMPIÓN\PRESENTACIONES CAPACITACIÓN\Vacunacion_SRP_SR_Presentacion_.pptx"
    create_improved_presentation(json_path, output_path)
