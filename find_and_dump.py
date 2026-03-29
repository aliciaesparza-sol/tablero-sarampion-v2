from pptx import Presentation
import sys
import io

# Ensure UTF-8 output for terminal
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
import sys

file_path = r'c:\Users\aicil\OneDrive\Escritorio\PVU\COEVA\PRESENTACIONES\CAPACITACIÓN SARAMPIÓN CLÍNICA Y VACUNACIÓN 17FEBRERO 2026.pptx'
prs = Presentation(file_path)

print("--- SEARCH RESULTS ---")
for i, slide in enumerate(prs.slides):
    text = ""
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            text += shape.text.upper() + " | "
    
    if "OPERACIONAL" in text or "DEFINICION" in text or "DEFINICIÓN" in text:
        print(f"SLIDE {i} (Possible Operational Definition):")
        print(text)
        print("-" * 20)
    
    if "DIFERENCIAL" in text:
        print(f"SLIDE {i} (Diagnostic Differential):")
        print(text)
        print("-" * 20)
