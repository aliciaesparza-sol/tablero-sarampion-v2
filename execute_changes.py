from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

file_path = r'c:\Users\aicil\OneDrive\Escritorio\PVU\COEVA\PRESENTACIONES\CAPACITACIÓN SARAMPIÓN CLÍNICA Y VACUNACIÓN 17FEBRERO 2026.pptx'
prs = Presentation(file_path)

# --- Change 1: Add "Fuente: OPS" to Slide 9 ---
slide_9 = prs.slides[9]
# Add a small text box at the bottom right
left = Inches(7)
top = Inches(6.8)
width = Inches(2.5)
height = Inches(0.5)
txBox = slide_9.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "Fuente: OPS"
p.font.size = Pt(12)
p.alignment = PP_ALIGN.RIGHT

# --- Change 2: Add Diagnostic Table to Slide 18 ---
slide_18 = prs.slides[18]

# Data for the table
headers = ["Enfermedad", "Etiología", "Hallazgos clínicos", "Características del exantema", "Laboratorio", "Clave para diferenciar"]
data = [
    ["Sarampión", "Paramyxovirus", "Fiebre alta (>38.5 °C), tos, coriza, conjuntivitis, manchas de Koplik (patognomónicas)", "Maculopapular, inicia retroauricular y línea capilar, progresión cefalocaudal, confluente", "IgM sarampión +, PCR +, leucopenia, linfopenia", "Tríada clásica + Koplik. Fiebre persiste 2–4 días tras exantema. Complicaciones respiratorias frecuentes."],
    ["Dengue", "Flavivirus (DENV 1–4)", "Fiebre alta bifásica, cefalea intensa, dolor retroorbitario, mialgias, artralgias (\"huesos rotos\")", "Petequial o maculopapular eritematoso, aparece 3–5 días post-fiebre; \"islas blancas en mar rojo\"", "NS1 Ag +, IgM/IgG dengue +, PCR +; trombocitopenia, hemoconcentración, leucopenia", "Dolor retroorbitario + trombocitopenia. Sin manchas de Koplik. Epidemiología vectorial."],
    ["Enfermedad de Kawasaki", "Vasculitis sistémica (etiología desconocida)", "Fiebre >5 días, conjuntivitis bilateral no purulenta, labios eritematosos/fisurados, lengua en fresa, adenopatía cervical, edema/eritema de extremidades", "Polimorfo (maculopapular, morbiliforme, escarlatiniforme); descamación periungueal tardía; eritema perianal", "Leucocitosis, PCR y VSG muy elevadas, trombocitosis en fase subaguda, ecocardiograma (aneurismas coronarios)", "5 criterios diagnósticos + fiebre >5 días. Sin agente infeccioso. Riesgo de aneurismas coronarios. Responde a IGIV + AAS."],
    ["Adenovirus", "Adenoviridae", "Fiebre alta, faringitis exudativa, conjuntivitis folicular (síndrome faringo-conjuntival), tos, diarrea en niños pequeños", "Maculopapular leve y transitorio, inespecífico; menos prominente que sarampión; puede estar ausente", "PCR adenovirus + (nasofaringe o heces), cultivo viral, leucocitosis moderada", "Síndrome faringo-conjuntival. Exantema inconstante y discreto. Sin manchas de Koplik. Brotes en colectividades."],
    ["Rubéola", "Rubivirus", "Fiebre leve o ausente, adenopatías retroauriculares y suboccipitales prominentes, artralgias en adultos", "Maculopapular fino, rosado, no confluente, inicio facial, duración ≈ 3 días", "IgM rubéola +, PCR rubéola +, hemograma habitualmente normal", "Adenopatías retroauriculares + cuadro leve. Sin tos ni Koplik. Riesgo teratogénico en embarazo."],
    ["Exantema súbito (roséola)", "HHV-6 / HHV-7", "Fiebre alta brusca 3–5 días en menores de 2 años; exantema aparece al ceder la fiebre", "Máculo-papular rosado, tronco y cuello, no confluente, dura 1–2 días", "Leucocitosis inicial, luego linfocitosis; IgM HHV-6 + (puede ser útil)", "Fiebre desaparece → exantema aparece. Exclusivo de lactantes. Sin fase catarral. Diagnóstico clínico."],
    ["Eritema infeccioso", "Parvovirus B19", "Fiebre leve o ausente, artralgias (más en adultos), crisis aplásica en pacientes con hemoglobinopatías", "\"Mejillas abofeteadas\", luego patrón reticular en encaje en extremidades; recurre con calor o ejercicio", "IgM Parvovirus B19 +, PCR +; anemia aplásica transitoria en inmunodeprimidos", "Cara en \"mejillas abofeteadas\" + exantema reticular. Sin fase catarral. Riesgo en embarazo e inmunodeprimidos."],
    ["Escarlatina", "S. pyogenes (EBHGA)", "Fiebre, odinofagia intensa, faringitis exudativa, lengua en fresa, palidez perioral", "Puntiforme, áspero al tacto (\"papel de lija\"), rubor difuso, pliegues de Pastia en axilas e ingles", "Exudado faríngeo + S. pyogenes, TASO elevado, leucocitosis con neutrofilia", "Faringitis bacteriana + textura áspera al tacto. Pliegues de Pastia patognomónicos. Responde a penicilina."]
]

rows = len(data) + 1
cols = len(headers)
# Positioning the table
left = Inches(0.5)
top = Inches(1.2)
width = Inches(12.3) # Assuming a 16:9 or 4:3 slide, wide table
height = Inches(5.5)

shape = slide_18.shapes.add_table(rows, cols, left, top, width, height)
table = shape.table

# Set header row
for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    # Style header
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(9)
            run.font.bold = True

# Fill data rows
for r_idx, row_data in enumerate(data):
    for c_idx, cell_data in enumerate(row_data):
        cell = table.cell(r_idx + 1, c_idx)
        cell.text = cell_data
        # Style cell
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(8)

prs.save(file_path)
print("Presentation modified successfully.")
