import pandas as pd
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import os
import io

# =================== DATA ===================
stock_file = r"c:\Users\aicil\OneDrive\Escritorio\PVU\SARAMPIÓN\EXISTENCIAS SRP Y SR\EXISTENCIAS_SRP_SR_POR_DIA.xlsx"
df_stock_raw = pd.read_excel(stock_file, sheet_name='Detalle por Jurisdicción')
df_stock_raw['FECHA_PARSED'] = pd.to_datetime(df_stock_raw['FECHA'], dayfirst=True, errors='coerce')
df_stock_raw['SECCIÓN'] = df_stock_raw['SECCIÓN'].astype(str).str.strip()

def normalize_jur(s):
    s = s.upper().strip()
    if 'NO. 1' in s or 'NO.1' in s: return 'JUR1'
    if 'NO. 2' in s or 'NO.2' in s: return 'JUR2'
    if 'NO. 3' in s or 'NO.3' in s: return 'JUR3'
    if 'NO. 4' in s or 'NO.4' in s: return 'JUR4'
    return None

df_stock_raw['JUR'] = df_stock_raw['SECCIÓN'].apply(normalize_jur)
df_jur = df_stock_raw[df_stock_raw['JUR'].notna() & df_stock_raw['FECHA_PARSED'].notna()].copy()
df_jur = df_jur.sort_values(['JUR', 'FECHA_PARSED'])
df_jur['CONSUMO_SR'] = df_jur.groupby('JUR')['SR PUNTOS'].diff().mul(-1).fillna(0)
df_jur['CONSUMO_SRP'] = df_jur.groupby('JUR')['SRP PUNTOS'].diff().mul(-1).fillna(0)

csv_path = r"c:\Users\aicil\OneDrive\Escritorio\PVU\SARAMPIÓN\REPORTE_SRP-SR-CENSIA\SRP-SR-2025_14-03-2026 06-41-16.csv"
df_doses = pd.read_csv(csv_path, usecols=['JURISDICCION','Fecha de registro','SR PRIMERA TOTAL','SR SEGUNDA TOTAL','SRP  PRIMERA TOTAL','SRP SEGUNDA TOTAL'])
df_doses['FECHA_PARSED'] = pd.to_datetime(df_doses['Fecha de registro'], errors='coerce')
df_doses = df_doses[df_doses['FECHA_PARSED'].dt.year >= 2026]
df_doses['SR_APLICADAS'] = df_doses[['SR PRIMERA TOTAL','SR SEGUNDA TOTAL']].fillna(0).sum(axis=1)
df_doses['SRP_APLICADAS'] = df_doses[['SRP  PRIMERA TOTAL','SRP SEGUNDA TOTAL']].fillna(0).sum(axis=1)

jur_map = {'DURANGO':'JUR1','SANTIAGO PAPASQUIARO':'JUR2','RODEO':'JUR3','MEZQUITAL':'JUR4'}
def map_jur(j):
    j = str(j).strip().upper()
    for k, v in jur_map.items():
        if k in j: return v
    return j

df_doses['JUR'] = df_doses['JURISDICCION'].apply(map_jur)
doses_agg = df_doses.groupby(['FECHA_PARSED','JUR']).agg(SR_APLICADAS=('SR_APLICADAS','sum'),SRP_APLICADAS=('SRP_APLICADAS','sum')).reset_index()

merged = pd.merge(df_jur[['FECHA_PARSED','JUR','SR PUNTOS','SRP PUNTOS','CONSUMO_SR','CONSUMO_SRP']],
                  doses_agg, on=['FECHA_PARSED','JUR'], how='left').fillna(0)
merged['DIFER_SR'] = merged['CONSUMO_SR'] - merged['SR_APLICADAS']
merged['DIFER_SRP'] = merged['CONSUMO_SRP'] - merged['SRP_APLICADAS']
merged = merged.sort_values(['FECHA_PARSED','JUR'])
merged['FECHA_STR'] = merged['FECHA_PARSED'].dt.strftime('%d/%m')

JUR_LABELS = {'JUR1':'Jur. 1 (Durango)','JUR2':'Jur. 2 (Santiago Papasquiaro)','JUR3':'Jur. 3 (Rodeo)','JUR4':'Jur. 4 (Mezquital)'}
JUR_FULL   = {'JUR1':'Jurisdicción No. 1 - Durango','JUR2':'Jurisdicción No. 2 - Santiago Papasquiaro','JUR3':'Jurisdicción No. 3 - Rodeo','JUR4':'Jurisdicción No. 4 - Mezquital'}
COLORS = {'JUR1':'#1F3864','JUR2':'#2196F3','JUR3':'#388E3C','JUR4':'#F57C00'}
DARK_BLUE = RGBColor(0x1F, 0x38, 0x64)
LIGHT_BLUE = RGBColor(0xBD, 0xD7, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

IMGS = {}
def fig_to_stream():
    buf = io.BytesIO()
    plt.savefig(buf, format='png', bbox_inches='tight', dpi=150)
    buf.seek(0)
    plt.close('all')
    return buf

# =================== CHARTS ===================
def make_sr_stock_chart():
    fig, ax = plt.subplots(figsize=(10, 4.5))
    for jur, jlabel in JUR_LABELS.items():
        sub = merged[merged['JUR']==jur].sort_values('FECHA_PARSED')
        if sub.empty: continue
        ax.plot(range(len(sub)), sub['SR PUNTOS'], marker='o', markersize=3, label=jlabel, color=COLORS[jur], linewidth=2)
    ticks = merged[merged['JUR']=='JUR1'].reset_index()
    ax.set_xticks(range(len(ticks)))
    ax.set_xticklabels(ticks['FECHA_STR'], rotation=45, ha='right', fontsize=7)
    ax.set_title('Existencias SR en Puntos de Vacunación por Jurisdicción', fontweight='bold', fontsize=12)
    ax.set_ylabel('Dosis en existencia')
    ax.grid(axis='y', alpha=0.3)
    ax.legend(fontsize=8)
    plt.tight_layout()
    return fig_to_stream()

def make_srp_stock_chart():
    fig, ax = plt.subplots(figsize=(10, 4.5))
    for jur, jlabel in JUR_LABELS.items():
        sub = merged[merged['JUR']==jur].sort_values('FECHA_PARSED')
        if sub.empty: continue
        ax.plot(range(len(sub)), sub['SRP PUNTOS'], marker='o', markersize=3, label=jlabel, color=COLORS[jur], linewidth=2)
    ticks = merged[merged['JUR']=='JUR1'].reset_index()
    ax.set_xticks(range(len(ticks)))
    ax.set_xticklabels(ticks['FECHA_STR'], rotation=45, ha='right', fontsize=7)
    ax.set_title('Existencias SRP en Puntos de Vacunación por Jurisdicción', fontweight='bold', fontsize=12)
    ax.set_ylabel('Dosis en existencia')
    ax.grid(axis='y', alpha=0.3)
    ax.legend(fontsize=8)
    plt.tight_layout()
    return fig_to_stream()

def make_aplicadas_chart(vaccine='SR'):
    col_ap = 'SR_APLICADAS' if vaccine=='SR' else 'SRP_APLICADAS'
    fig, axes = plt.subplots(2, 2, figsize=(12, 7), sharey=False)
    axes = axes.flatten()
    for i, (jur, jlabel) in enumerate(JUR_LABELS.items()):
        sub = merged[merged['JUR']==jur].sort_values('FECHA_PARSED')
        ax = axes[i]
        colors = ['#C00000' if v > 0 else '#BDD7EE' for v in sub[col_ap]]
        ax.bar(range(len(sub)), sub[col_ap], color=colors, width=0.8)
        ax.set_title(jlabel, fontweight='bold', fontsize=8)
        ax.set_xticks(range(len(sub)))
        ax.set_xticklabels(sub['FECHA_STR'], rotation=60, fontsize=5)
        ax.grid(axis='y', alpha=0.3)
        ax.set_ylabel('Dosis', fontsize=7)
    fig.suptitle(f'Dosis {vaccine} Aplicadas por Jurisdicción y Día', fontweight='bold', fontsize=13)
    plt.tight_layout()
    return fig_to_stream()

def make_consumo_chart(vaccine='SR'):
    col_c = 'CONSUMO_SR' if vaccine=='SR' else 'CONSUMO_SRP'
    fig, axes = plt.subplots(2, 2, figsize=(12, 7), sharey=False)
    axes = axes.flatten()
    for i, (jur, jlabel) in enumerate(JUR_LABELS.items()):
        sub = merged[merged['JUR']==jur].sort_values('FECHA_PARSED')
        ax = axes[i]
        colors = ['#1F3864' if v >= 0 else '#FF7043' for v in sub[col_c]]
        ax.bar(range(len(sub)), sub[col_c], color=colors, width=0.8)
        ax.set_title(jlabel, fontweight='bold', fontsize=8)
        ax.set_xticks(range(len(sub)))
        ax.set_xticklabels(sub['FECHA_STR'], rotation=60, fontsize=5)
        ax.axhline(0, color='black', linewidth=0.7)
        ax.grid(axis='y', alpha=0.3)
        ax.set_ylabel('Δ Dosis', fontsize=7)
    fig.suptitle(f'Consumo Teórico (Δ Stock) {vaccine} por Jurisdicción y Día', fontweight='bold', fontsize=13)
    plt.tight_layout()
    return fig_to_stream()

def make_corr_chart(vaccine='SR'):
    col_ap = 'SR_APLICADAS' if vaccine=='SR' else 'SRP_APLICADAS'
    col_c  = 'CONSUMO_SR'   if vaccine=='SR' else 'CONSUMO_SRP'
    fig, axes = plt.subplots(2, 2, figsize=(12, 7), sharey=False)
    axes = axes.flatten()
    for i, (jur, jlabel) in enumerate(JUR_LABELS.items()):
        sub = merged[merged['JUR']==jur].sort_values('FECHA_PARSED')
        ax = axes[i]
        x = range(len(sub))
        ax.bar(x, sub[col_c],   width=0.6, label='Consumo Teórico', color='#1F3864', alpha=0.8)
        ax.bar(x, sub[col_ap],  width=0.3, label='Dosis Aplicadas', color='#FF7043', alpha=0.9)
        ax.set_title(jlabel, fontweight='bold', fontsize=8)
        ax.set_xticks(list(x))
        ax.set_xticklabels(list(sub['FECHA_STR']), rotation=60, fontsize=5)
        ax.axhline(0, color='black', linewidth=0.7)
        ax.grid(axis='y', alpha=0.3)
        ax.set_ylabel('Dosis', fontsize=7)
        if i == 0: ax.legend(fontsize=7)
    fig.suptitle(f'Correlación: Consumo Teórico vs Dosis Aplicadas ({vaccine})', fontweight='bold', fontsize=13)
    plt.tight_layout()
    return fig_to_stream()

# Generate all charts
charts = {
    'sr_stock':   make_sr_stock_chart(),
    'srp_stock':  make_srp_stock_chart(),
    'sr_aplicadas': make_aplicadas_chart('SR'),
    'srp_aplicadas': make_aplicadas_chart('SRP'),
    'sr_consumo': make_consumo_chart('SR'),
    'srp_consumo': make_consumo_chart('SRP'),
    'sr_corr':  make_corr_chart('SR'),
    'srp_corr': make_corr_chart('SRP')
}

# =================== PPTX ===================
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

def add_slide(title_text, subtitle_text=None, bg_dark=False):
    slide = prs.slides.add_slide(blank)
    bg = slide.background.fill
    if bg_dark:
        bg.solid()
        bg.fore_color.rgb = DARK_BLUE
    return slide

def add_textbox(slide, text, left, top, width, height, size=14, bold=False, color=RGBColor(0x1F,0x38,0x64), align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def add_header_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()
    add_textbox(slide, title, 0.2, 0.05, 12.5, 0.7, size=24, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle, 0.2, 0.72, 12.5, 0.5, size=13, color=LIGHT_BLUE, align=PP_ALIGN.LEFT)

def add_image_stream(slide, stream, left, top, width, height):
    slide.shapes.add_picture(stream, Inches(left), Inches(top), Inches(width), Inches(height))

def add_bullet_box(slide, bullets, left, top, width, height, title=None, bg=None):
    if bg:
        box = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
        box.fill.solid()
        box.fill.fore_color.rgb = bg
        box.line.color.rgb = RGBColor(0xBD,0xD7,0xEE)
    tb = slide.shapes.add_textbox(Inches(left+0.1), Inches(top+0.05), Inches(width-0.2), Inches(height-0.1))
    tf = tb.text_frame
    tf.word_wrap = True
    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.bold = True
        run.font.size = Pt(12)
        run.font.color.rgb = DARK_BLUE
    for b in bullets:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"▸  {b}"
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(0x20,0x20,0x20)

# === Slide 1: Portada ===
slide = add_slide("", bg_dark=True)
add_textbox(slide, "SECRETARÍA DE SALUD — DURANGO", 1, 1.2, 11, 0.6, size=14, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_textbox(slide, "Análisis de Existencias, Consumo\ny Dosis Aplicadas de Vacunas SR y SRP", 0.5, 2.0, 12.3, 2.5, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "Campaña de Vacunación contra Sarampión 2026", 1, 4.3, 11, 0.6, size=16, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_textbox(slide, "Periodo: 18 de febrero — 14 de marzo de 2026", 1, 4.9, 11, 0.5, size=13, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_textbox(slide, "Presentado ante el COEVA — Marzo 2026", 1, 5.5, 11, 0.5, size=12, color=RGBColor(0x90,0xCA,0xF9), align=PP_ALIGN.CENTER)

# === Slide 2: Contexto ===
slide = add_slide("")
add_header_bar(slide, "1. Contexto y Metodología", "¿Qué se analizó y cómo?")
add_bullet_box(slide, [
    "Periodo analizado: 18 de febrero al 14 de marzo de 2026 (25 días hábiles)",
    "Fuentes de información: ① Registro diario de existencias por jurisdicción (EXISTENCIAS_SRP_SR_POR_DIA.xlsx) ② Base de datos CENSIA de dosis aplicadas (SRP-SR-2025_14-03-2026.csv)",
    "Se analizaron 4 jurisdicciones sanitarias del estado de Durango",
    "Variables estudiadas:\n   — Existencias en puntos de vacunación (SR y SRP)\n   — Consumo teórico = diferencia de stock entre días consecutivos\n   — Dosis aplicadas registradas en el sistema CENSIA",
    "Se calculó la diferencia entre el consumo teórico y las dosis aplicadas para detectar posibles inconsistencias en el registro",
], 0.3, 1.4, 12.7, 5.8, bg=RGBColor(0xF0,0xF5,0xFF))

# === Slide 3: Stock SR ===
slide = add_slide("")
add_header_bar(slide, "2. Existencias SR por Jurisdicción", "Stock de vacuna SR en puntos de vacunación — 18/02 al 14/03/2026")
add_image_stream(slide, charts['sr_stock'], 0.2, 1.35, 12.9, 6.0)

# === Slide 4: Stock SRP ===
slide = add_slide("")
add_header_bar(slide, "3. Existencias SRP por Jurisdicción", "Stock de vacuna SRP en puntos de vacunación — 18/02 al 14/03/2026")
add_image_stream(slide, charts['srp_stock'], 0.2, 1.35, 12.9, 6.0)

# === Slide 5: Dosis aplicadas SR ===
slide = add_slide("")
add_header_bar(slide, "4. Dosis SR Aplicadas por Jurisdicción y Día", "Registro de aplicación de vacuna SR (Doble viral) por jurisdicción")
add_image_stream(slide, charts['sr_aplicadas'], 0.2, 1.35, 12.9, 6.0)

# === Slide 6: Dosis aplicadas SRP ===
slide = add_slide("")
add_header_bar(slide, "5. Dosis SRP Aplicadas por Jurisdicción y Día", "Registro de aplicación de vacuna SRP (Triple viral) por jurisdicción")
add_image_stream(slide, charts['srp_aplicadas'], 0.2, 1.35, 12.9, 6.0)

# === Slide 7: Consumo SR ===
slide = add_slide("")
add_header_bar(slide, "6. Consumo Teórico SR (Δ Stock)", "Variación neta del stock de SR entre días consecutivos — positivo = consumo real, negativo = reabastecimiento")
add_image_stream(slide, charts['sr_consumo'], 0.2, 1.35, 12.9, 6.0)

# === Slide 8: Consumo SRP ===
slide = add_slide("")
add_header_bar(slide, "7. Consumo Teórico SRP (Δ Stock)", "Variación neta del stock de SRP entre días consecutivos")
add_image_stream(slide, charts['srp_consumo'], 0.2, 1.35, 12.9, 6.0)

# === Slide 9: Correlacion SR ===
slide = add_slide("")
add_header_bar(slide, "8. Correlación: Consumo Teórico vs Dosis Aplicadas — SR", "Comparación entre la variación en stock y el registro de aplicación en CENSIA")
add_image_stream(slide, charts['sr_corr'], 0.2, 1.35, 12.9, 6.0)

# === Slide 10: Correlacion SRP ===
slide = add_slide("")
add_header_bar(slide, "9. Correlación: Consumo Teórico vs Dosis Aplicadas — SRP", "Comparación entre la variación en stock y el registro de aplicación en CENSIA")
add_image_stream(slide, charts['srp_corr'], 0.2, 1.35, 12.9, 6.0)

# === Slide 11: Hallazgos por Jurisdicción ===
findings = {
    'JUR1': [
        "Stock SR inicial: 350 → final: 12,090 dosis (+11,740 neto con reabastecimientos)",
        "Stock SRP inicial: 5,127 → final: 415 dosis (−4,712 consumidas en el periodo)",
        "Pico de actividad: 4-6 marzo con stocks SR mayores a 13,000 dosis distribuidas a puntos",
        "Se detectaron movimientos de stock sin registro de aplicación en fechas 03–06 y 09–11 marzo",
    ],
    'JUR2': [
        "Stock SR inicial: 0 → final: 5,190 dosis (recibió vacunas hacia final del periodo)",
        "Stock SRP inicial: 2,014 → final: 989 dosis (consumo acumulado: ~1,025 dosis SRP)",
        "Los días 10–13 marzo muestran desfases mayores a 300 dosis SR entre stock y registro CENSIA",
        "Presenta una de las mayores brechas acumuladas en SR entre jurisdicciones",
    ],
    'JUR3': [
        "Jurisdicción con menor volumen de vacunas: SR stock entre 60–390 dosis en todo el periodo",
        "SRP: consumo acumulado de ~911 dosis, con stock reduciéndose de 2,211 a 1,300",
        "Bajo nivel de reportes de aplicación en CENSIA para esta jurisdicción",
        "Requiere refuerzo en el registro diario de dosis aplicadas",
    ],
    'JUR4': [
        "Stock SR oscila entre 140–1,072 dosis en el periodo analizado",
        "Consumo teórico SR acumulado: 712 dosis; SRP: 235 dosis",
        "Presenta movimientos de stock el 02/03 y 09/03 sin registro en CENSIA",
        "Desfases medios moderados (SR: ~168 dosis/día, SRP: ~99 dosis/día)",
    ]
}

for jur, bullets in findings.items():
    slide = add_slide("")
    add_header_bar(slide, f"10. Hallazgos — {JUR_FULL[jur]}", "Análisis de existencias, consumo y dosis aplicadas")
    add_bullet_box(slide, bullets, 0.3, 1.4, 12.7, 5.6, bg=RGBColor(0xF0,0xF5,0xFF))

# === Slide 15: Conclusiones ===
slide = add_slide("")
add_header_bar(slide, "11. Conclusiones y Recomendaciones", "")
add_bullet_box(slide, [
    "✅  Los stocks de biológicos SR y SRP se distribuyeron oportunamente a los 4 jurisdicciones a lo largo del periodo.",
    "⚠️  Se identificaron desfases entre el consumo teórico (variación de stock) y las dosis registradas en CENSIA en múltiples jurisdicciones, especialmente durante el periodo 03–11 de marzo.",
    "⚠️  La Jurisdicción No. 1 (Durango) presenta el mayor número de días con desfases altos en SR, lo que podría indicar subreporte de dosis aplicadas en el sistema.",
    "⚠️  La Jurisdicción No. 3 (Rodeo) reporta muy pocas dosis aplicadas en CENSIA a pesar de contar con existencias distribuidas.",
    "📋  Recomendación 1: Depurar y actualizar el registro de dosis aplicadas en CENSIA para los días con desfases identificados.",
    "📋  Recomendación 2: Establecer un proceso de conciliación diaria entre el reporte de existencias y el registro de aplicación.",
    "📋  Recomendación 3: Verificar que los puntos de vacunación de las 4 jurisdicciones reporten diariamente al sistema.",
    "📋  Recomendación 4: Revisar posibles ingresos de stock no programados que expliquen los aumentos en existencias en fechas sin aplicación reportada.",
], 0.3, 1.4, 12.7, 5.8, bg=RGBColor(0xF0,0xF5,0xFF))

out_path = r"c:\Users\aicil\OneDrive\Escritorio\PVU\SARAMPIÓN\ANALISIS_EXISTENCIAS_COEVA_2026.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
