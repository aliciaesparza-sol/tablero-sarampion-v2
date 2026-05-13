import os
import urllib.request
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def download_image(url, filename):
    try:
        req = urllib.request.Request(url, headers={'User-Agent': 'Mozilla/5.0'})
        with urllib.request.urlopen(req) as response, open(filename, 'wb') as out_file:
            data = response.read()
            out_file.write(data)
        return True
    except Exception as e:
        print(f"Failed to download {url}: {e}")
        return False

def update_presentation(pptx_path):
    prs = Presentation(pptx_path)
    
    # Download images
    img1 = "vaccine.jpg"
    img2 = "measles.jpg"
    download_image("https://upload.wikimedia.org/wikipedia/commons/thumb/1/12/Syringe_and_vaccine.jpg/800px-Syringe_and_vaccine.jpg", img1)
    download_image("https://upload.wikimedia.org/wikipedia/commons/thumb/2/25/Measles_child.jpg/800px-Measles_child.jpg", img2)
    
    # Slide 3 (Index 2): Tipo de vacuna. Add vaccine image
    if len(prs.slides) > 2:
        slide3 = prs.slides[2]
        if os.path.exists(img1):
            slide3.shapes.add_picture(img1, Inches(7.5), Inches(2), width=Inches(4))

    # Slide 9 (Index 8): ESAVI / Sarampion. Add measles image
    if len(prs.slides) > 8:
        slide9 = prs.slides[8]
        if os.path.exists(img2):
            slide9.shapes.add_picture(img2, Inches(7.5), Inches(2), width=Inches(4))

    # Add new slide for missing info (Lineamientos Brote 2025-2026)
    # Use the same layout as the second slide (usually Title and Content)
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    new_slide = prs.slides.add_slide(layout)
    
    if new_slide.shapes.title:
        new_slide.shapes.title.text = "Actualización Lineamientos Nacionales 2025-2026 (Brote)"
        
    for shape in new_slide.shapes:
        if shape.has_text_frame and shape != new_slide.shapes.title:
            tf = shape.text_frame
            tf.text = "Estrategia Intensiva:"
            
            p1 = tf.add_paragraph()
            p1.text = "Dosis Cero: Aplicación en menores de 6 a 11 meses (no sustituye la de 12 meses)."
            p1.level = 1
            
            p2 = tf.add_paragraph()
            p2.text = "Población 13-49 años: Vacunación con SR a susceptibles."
            p2.level = 1
            
            p3 = tf.add_paragraph()
            p3.text = "Personal de Salud: Refuerzo obligatorio (1 dosis adicional)."
            p3.level = 1
            
            p4 = tf.add_paragraph()
            p4.text = "Bloqueo Vacunal: Vacunación casa por casa ante casos confirmados."
            p4.level = 1
            
            p5 = tf.add_paragraph()
            p5.text = "Validación documental: Requiere comprobación en Cartilla Nacional de Salud."
            p5.level = 1
            break
            
    try:
        prs.save(pptx_path)
        print(f"Successfully updated {pptx_path}")
    except Exception as e:
        print(f"Error saving: {e}")

if __name__ == '__main__':
    target_path = r"c:\Users\aicil\OneDrive\Escritorio\PVU\SARAMPIÓN\PRESENTACIONES CAPACITACIÓN\Vacunacion_SRP_SR_Presentacion_.pptx"
    update_presentation(target_path)
