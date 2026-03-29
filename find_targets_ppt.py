from pptx import Presentation

file_path = r'c:\Users\aicil\OneDrive\Escritorio\PVU\COEVA\PRESENTACIONES\CAPACITACIÓN SARAMPIÓN CLÍNICA Y VACUNACIÓN 17FEBRERO 2026.pptx'
search_terms = ["DEFINICIÓN", "OPERACIONAL", "DIFERENCIAL", "SOSPECHOSO", "ALGORITMO"]

try:
    prs = Presentation(file_path)
    for i, slide in enumerate(prs.slides):
        found = False
        text_content = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    t = p.text.upper()
                    if any(term in t for term in search_terms):
                        found = True
                    text_content.append(p.text)
        if found:
            print(f"\n--- SLIDE {i} ---")
            for t in text_content:
                if t.strip():
                    print(t)
except Exception as e:
    print(f"Error: {e}")
