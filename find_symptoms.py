from pptx import Presentation
import re

file_path = r'c:\Users\aicil\OneDrive\Escritorio\PVU\COEVA\PRESENTACIONES\CAPACITACIÓN SARAMPIÓN CLÍNICA Y VACUNACIÓN 17FEBRERO 2026.pptx'
prs = Presentation(file_path)

for i, slide in enumerate(prs.slides):
    text = ""
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            text += shape.text.upper() + " "
    
    if "FIEBRE" in text and "EXANTEMA" in text:
        print(f"Slide {i}: Contains 'FIEBRE' and 'EXANTEMA'")
        print(f"Content: {text[:200]}...")
