import collections.abc
from pptx import Presentation
import os
import json

def extract_pptx_data(pptx_path, output_dir):
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    prs = Presentation(pptx_path)
    data = []
    
    image_count = 0
    for i, slide in enumerate(prs.slides):
        slide_data = {
            "slide_index": i + 1,
            "title": "",
            "text": [],
            "images": [],
            "notes": ""
        }
        
        # Extract title
        if slide.shapes.title:
            slide_data["title"] = slide.shapes.title.text
            
        # Extract text from all shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                if shape != slide.shapes.title:
                    slide_data["text"].append(shape.text.strip())
            
            # Extract images
            if shape.shape_type == 13: # Picture
                image_count += 1
                image = shape.image
                image_filename = f"slide_{i+1}_img_{image_count}.{image.ext}"
                image_path = os.path.join(output_dir, image_filename)
                with open(image_path, "wb") as f:
                    f.write(image.blob)
                slide_data["images"].append(image_filename)
        
        # Extract notes
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            slide_data["notes"] = notes
            
        data.append(slide_data)
        
    return data

pptx_file = r"C:\Users\aicil\.gemini\antigravity\scratch\sarampion_tmp.pptx"
output_folder = r"C:\Users\aicil\.gemini\antigravity\scratch\measles_presentation\assets"
json_output = r"C:\Users\aicil\.gemini\antigravity\scratch\detailed_ppt_data.json"

try:
    results = extract_pptx_data(pptx_file, output_folder)
    with open(json_output, "w", encoding="utf-8") as f:
        json.dump(results, f, ensure_ascii=False, indent=4)
    print(f"Extraction successful. Data saved to {json_output}")
except Exception as e:
    print(f"Error: {e}")
