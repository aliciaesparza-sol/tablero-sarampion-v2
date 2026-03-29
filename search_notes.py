from pptx import Presentation

file_path = r'c:\Users\aicil\OneDrive\Escritorio\PVU\COEVA\PRESENTACIONES\CAPACITACIÓN SARAMPIÓN CLÍNICA Y VACUNACIÓN 17FEBRERO 2026.pptx'
prs = Presentation(file_path)

for i, slide in enumerate(prs.slides):
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text
        if "DEFINIC" in notes.upper() or "OPERACIONAL" in notes.upper():
            print(f"Match Slide {i} NOTES: {notes}")
