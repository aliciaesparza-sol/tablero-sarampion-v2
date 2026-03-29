from pptx import Presentation

file_path = r'c:\Users\aicil\OneDrive\Escritorio\PVU\COEVA\PRESENTACIONES\CAPACITACIÓN SARAMPIÓN CLÍNICA Y VACUNACIÓN 17FEBRERO 2026.pptx'
prs = Presentation(file_path)

for i, slide in enumerate(prs.slides):
    text = ""
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            text += shape.text + " "
    
    if "DEFINICIÓN" in text.upper() or "DEFINICION" in text.upper() or "OPERACIONAL" in text.upper():
        print(f"Match Slide {i}: {text[:100].replace('\n', ' ')}...")
    if "DIAGNOSTICO" in text.upper() or "DIAGNÓSTICO" in text.upper() or "DIFERENCIAL" in text.upper():
        print(f"Match Slide {i}: {text[:100].replace('\n', ' ')}...")
