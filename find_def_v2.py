from pptx import Presentation
import re

file_path = r'c:\Users\aicil\OneDrive\Escritorio\PVU\COEVA\PRESENTACIONES\CAPACITACIÓN SARAMPIÓN CLÍNICA Y VACUNACIÓN 17FEBRERO 2026.pptx'
prs = Presentation(file_path)

keywords = ["FEBRIL", "EXANTEMÁTICA", "MACULOPAPULAR", "TOS", "CORIZA", "CONJUNTIVITIS", "CASO", "PROBABLE", "CONFIRMADO", "DESCARTADO"]

for i, slide in enumerate(prs.slides):
    text = ""
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text += shape.text.upper() + " "
    
    matches = [k for k in keywords if k in text]
    if len(matches) > 3:
        print(f"Slide {i}: Matches {matches}")
        print(f"Content: {text[:300]}...")
