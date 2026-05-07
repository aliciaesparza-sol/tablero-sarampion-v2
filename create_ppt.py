from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Paths to generated images
MAP_PATH = r"C:\Users\aicil\.gemini\antigravity\brain\b90ab7f9-3458-429d-99b9-0806d80e4e60\mapa_estrategia_aerea_mezquital_1778103476848.png"
BG_PATH = r"C:\Users\aicil\.gemini\antigravity\brain\b90ab7f9-3458-429d-99b9-0806d80e4e60\fondo_presentacion_salud_aerea_1778103527994.png"

def create_presentation():
    prs = Presentation()

    # Slide 1: Portada
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Add background
    if os.path.exists(BG_PATH):
        slide.shapes.add_picture(BG_PATH, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Title Box
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "ESTRATEGIA DE APOYO AÉREO"
    p.font.bold = True
    p.font.size = Pt(44)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Operativo Mezquital 2026 - Plan de Vuelos"
    p2.font.size = Pt(28)
    p2.font.color.rgb = RGBColor(200, 200, 200)
    p2.alignment = PP_ALIGN.CENTER

    # Slide 2: Introducción y Objetivos
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Objetivos de la Estrategia Aérea"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Garantizar el acceso a salud en zonas de difícil acceso."
    p = tf.add_paragraph()
    p.text = "Optimización logística para el despliegue de Unidades Médicas PFAM."
    p = tf.add_paragraph()
    p.text = "Coordinación efectiva entre Durango Capital y la Zona Indígena."

    # Slide 3: Mapa Estratégico
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Mapa de Operaciones: Mezquital"
    p.font.bold = True
    p.font.size = Pt(32)
    
    if os.path.exists(MAP_PATH):
        slide.shapes.add_picture(MAP_PATH, Inches(1), Inches(1.2), width=Inches(8))

    # Slide 4: Fase 1 - Despliegue (18 de Mayo)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Fase 1: Despliegue (18 de Mayo de 2026)"
    
    tf = slide.placeholders[1].text_frame
    tf.text = "08:00 AM: Salida de Durango hacia La Guajolota."
    p = tf.add_paragraph()
    p.text = "TRASLADO 1: Posicionamiento de Unidad Médica PFAM en Potreros (3 integrantes e insumos)."
    p = tf.add_paragraph()
    p.text = "TRASLADO 2: Posicionamiento de Unidad Médica PFAM en Cihuacora."
    p = tf.add_paragraph()
    p.text = "Retorno del equipo a Durango al finalizar el posicionamiento."

    # Slide 5: Permanencia y Actividad
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Operativo en Terreno (18-20 de Mayo)"
    
    tf = slide.placeholders[1].text_frame
    tf.text = "Las Unidades Médicas PFAM permanecen 3 días en Potreros y Cihuacora."
    p = tf.add_paragraph()
    p.text = "Actividades de salud, vacunación y atención integral."
    p = tf.add_paragraph()
    p.text = "Equipos ESI inician actividades paralelas en las zonas asignadas."

    # Slide 6: Fase 2 - Recuperación (20 de Mayo)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Fase 2: Recuperación (20 de Mayo de 2026)"
    
    tf = slide.placeholders[1].text_frame
    tf.text = "13:00 PM: Traslado de Durango hacia Cihuacora."
    p = tf.add_paragraph()
    p.text = "RETORNO 1: Traslado de Cihuacora a La Guajolota. ESI continúa en Llanos de Jacalitos."
    p = tf.add_paragraph()
    p.text = "RETORNO 2: Traslado de Potreros a La Guajolota. ESI continúa en Sombrero Quemado."
    p = tf.add_paragraph()
    p.text = "Finalización: Cierre del Plan de Vuelos y retorno a Durango Capital."

    # Slide 7: Conclusión
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Resumen de Impacto"
    
    tf = slide.placeholders[1].text_frame
    tf.text = "Cobertura garantizada en comunidades remotas."
    p = tf.add_paragraph()
    p.text = "Reducción de tiempos de traslado de 12 horas (terrestre) a minutos (aéreo)."
    p = tf.add_paragraph()
    p.text = "Eficiencia operativa y seguridad para el personal médico."

    output_path = r"C:\Users\aicil\.gemini\antigravity\scratch\Estrategia_Aerea_Mezquital_2026.pptx"
    prs.save(output_path)
    print(f"Presentación guardada en: {output_path}")

if __name__ == "__main__":
    create_presentation()
