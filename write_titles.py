from pptx import Presentation

file_path = r'c:\Users\aicil\OneDrive\Escritorio\PVU\COEVA\PRESENTACIONES\CAPACITACIÓN SARAMPIÓN CLÍNICA Y VACUNACIÓN 17FEBRERO 2026.pptx'
prs = Presentation(file_path)

with open('titles.txt', 'w', encoding='utf-8') as f:
    for i, slide in enumerate(prs.slides):
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text
        else:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    title = shape.text.strip().split('\n')[0][:50]
                    break
        f.write(f"Slide {i}: {title}\n")
