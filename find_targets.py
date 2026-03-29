from pptx import Presentation

file_path = r'c:\Users\aicil\OneDrive\Escritorio\PVU\COEVA\PRESENTACIONES\CAPACITACIÓN SARAMPIÓN CLÍNICA Y VACUNACIÓN 17FEBRERO 2026.pptx'
prs = Presentation(file_path)

print("Searching for Target Slides...")
for i, slide in enumerate(prs.slides):
    text = ""
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            text += shape.text.upper() + " "
    
    if "CASO PROBABLE" in text or "DEFINICIÓN OPERACIONAL" in text or "DEFINICION OPERACIONAL" in text:
        print(f"Slide {i}: Likely 'Definición Operacional' | Text: {text[:200]}...")
    
    if "DIFERENCIAL" in text:
        print(f"Slide {i}: Likely 'Diagnóstico Diferencial' | Text: {text[:200]}...")
