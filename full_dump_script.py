from pptx import Presentation

file_path = r'c:\Users\aicil\OneDrive\Escritorio\PVU\COEVA\PRESENTACIONES\CAPACITACIÓN SARAMPIÓN CLÍNICA Y VACUNACIÓN 17FEBRERO 2026.pptx'
prs = Presentation(file_path)

with open('full_dump.txt', 'w', encoding='utf-8') as f:
    for i, slide in enumerate(prs.slides):
        f.write(f"--- SLIDE {i} ---\n")
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                f.write(shape.text + "\n")
        f.write("\n")
