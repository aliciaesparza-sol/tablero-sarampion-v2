from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

file_path = r'c:\Users\aicil\OneDrive\Escritorio\PVU\COEVA\PRESENTACIONES\CAPACITACIÓN SARAMPIÓN CLÍNICA Y VACUNACIÓN 17FEBRERO 2026.pptx'
prs = Presentation(file_path)

for i, slide in enumerate(prs.slides):
    num_txt = sum(1 for s in slide.shapes if hasattr(s, "text") and s.text.strip())
    num_img = sum(1 for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
    print(f"Slide {i}: Text Shapes: {num_txt}, Images: {num_img}")
    if num_txt > 0:
        first_text = [s.text for s in slide.shapes if hasattr(s, "text") and s.text.strip()][0][:50]
        print(f"  First Text: {first_text}")
