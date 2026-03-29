from pptx import Presentation

file_path = r'c:\Users\aicil\OneDrive\Escritorio\PVU\COEVA\PRESENTACIONES\CAPACITACIÓN SARAMPIÓN CLÍNICA Y VACUNACIÓN 17FEBRERO 2026.pptx'
try:
    prs = Presentation(file_path)
    print(f"Total slides: {len(prs.slides)}")

    for i, slide in enumerate(prs.slides):
        print(f"\n--- SLIDE {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs)
                    if text.strip():
                        print(f"[{shape.name}] {text}")
            elif shape.has_table:
                print(f"[{shape.name}] (Table detected)")
                for row in shape.table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text_frame.text.strip())
                    print(f"  | {' | '.join(row_text)} |")
except Exception as e:
    print(f"Error: {e}")
