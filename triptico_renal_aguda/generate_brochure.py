import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_acute_pptx():
    prs = Presentation()
    prs.slide_width = Inches(11.69)
    prs.slide_height = Inches(8.27)

    def add_column_text(slide, x, y, width, title, bullets, title_size=24, bullet_size=15, color=(183, 28, 28)):
        # Title Background Header
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, Inches(0.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*color)
        shape.line.color.rgb = RGBColor(*color)
        
        txBox = slide.shapes.add_textbox(x, y, width, Inches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title.upper()
        p.font.bold = True
        p.font.size = Pt(title_size)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        curr_y = y + Inches(0.8)
        for bullet in bullets:
            txBox2 = slide.shapes.add_textbox(x + Inches(0.2), curr_y, width - Inches(0.4), Inches(0.4))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = f"• {bullet}"
            p2.font.size = Pt(bullet_size)
            p2.font.bold = True
            p2.font.color.rgb = RGBColor(0, 0, 0)
            curr_y += Inches(0.55)
        return curr_y

    # --- SLIDE 1: INTERIOR ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    col_width = prs.slide_width / 3
    
    # Column 1
    add_column_text(slide1, Inches(0.1), Inches(0.5), col_width - Inches(0.2), "¿QUÉ ES?", [
        "La Lesión Renal Aguda (LRA) es la disminución súbita de la función renal.",
        "Ocurre en horas o días y puede ser reversible si se trata a tiempo."
    ])
    add_column_text(slide1, Inches(0.1), Inches(3.5), col_width - Inches(0.2), "PRINCIPALES CAUSAS", [
        "Deshidratación", "Sepsis", "Choque", "Medicamentos nefrotóxicos"
    ])

    # Column 2
    add_column_text(slide1, col_width + Inches(0.1), Inches(0.5), col_width - Inches(0.2), "SIGNOS Y SÍNTOMAS", [
        "Oliguria o anuria", "Aumento rápido de peso", "Edema generalizado", 
        "Fatiga", "Confusión", "Alteraciones electrolíticas"
    ])
    
    # Column 3
    add_column_text(slide1, 2*col_width + Inches(0.1), Inches(0.5), col_width - Inches(0.2), "CUIDADOS DE ENFERMERÍA", [
        "Control estricto de signos vitales", "Balance hídrico horario", "Monitorización de diuresis", 
        "Administración segura de líquidos y fármacos", "Prevención de infecciones", "Vigilancia de electrolitos"
    ])
    if os.path.exists("nursing.png"):
        slide1.shapes.add_picture("nursing.png", 2*col_width + Inches(0.5), Inches(5.5), width=Inches(3))

    # --- SLIDE 2: EXTERIOR ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Column 1: EDUCACIÓN AL PACIENTE
    add_column_text(slide2, Inches(0.1), Inches(0.5), col_width - Inches(0.2), "EDUCACIÓN AL PACIENTE", [
        "No automedicarse", "Mantener hidratación adecuada", "Reportar cambios en la orina", 
        "Cumplir indicaciones médicas", "Asistir a seguimiento clínico"
    ])

    # Column 2: ¿CUÁNDO ACUDIR?
    add_column_text(slide2, col_width + Inches(0.1), Inches(0.5), col_width - Inches(0.2), "¿CUÁNDO ACUDIR?", [
        "Disminución o ausencia de orina", "Edema súbito", "Dificultad respiratoria", 
        "Confusión", "Náuseas o vómitos persistentes"
    ])
    
    label_box = slide2.shapes.add_textbox(col_width + Inches(0.1), Inches(7.5), col_width - Inches(0.2), Inches(0.5))
    label_box.text_frame.text = "SERVICIO DE ENFERMERÍA"
    p = label_box.text_frame.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(183, 28, 28)
    p.alignment = PP_ALIGN.CENTER

    # Column 3: PORTADA
    if os.path.exists("cover.png"):
        pic = slide2.shapes.add_picture("cover.png", 2*col_width, 0, width=col_width, height=prs.slide_height)
        # Add a semi-transparent overlay
        overlay = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 2*col_width, 0, col_width, prs.slide_height)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(183, 28, 28)
        overlay.fill.transparency = 0.2
        overlay.line.fill.background()

    cover_box = slide2.shapes.add_textbox(2*col_width + Inches(0.2), Inches(1), col_width - Inches(0.4), Inches(6))
    tf = cover_box.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "LESIÓN RENAL AGUDA"
    p.font.bold = True
    p.font.size = Pt(40)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "Cuidados de Enfermería"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(255, 200, 200)
    p2.alignment = PP_ALIGN.CENTER
    
    p3 = tf.add_paragraph()
    p3.text = "\nAtención oportuna • Vigilancia • Prevención"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(255, 255, 255)
    p3.alignment = PP_ALIGN.CENTER

    p_name = tf.add_paragraph()
    p_name.text = "\nALUMNA:\nMyriam Machado Arreola"
    p_name.font.size = Pt(18)
    p_name.font.color.rgb = RGBColor(255, 255, 255)
    p_name.font.bold = True
    p_name.alignment = PP_ALIGN.CENTER
    
    p_enf = tf.add_paragraph()
    p_enf.text = "\nSERVICIO DE ENFERMERÍA"
    p_enf.font.size = Pt(16)
    p_enf.font.bold = True
    p_enf.font.color.rgb = RGBColor(255, 255, 255)
    p_enf.alignment = PP_ALIGN.CENTER

    prs.save("Triptico_LRA_Final.pptx")

if __name__ == "__main__":
    create_acute_pptx()
