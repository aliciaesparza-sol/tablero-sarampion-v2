import sys; sys.path.append('.')
from pptx import Presentation

file_path = r'c:\Users\aicil\OneDrive\Escritorio\PVU\COEVA\PRESENTACIONES\CAPACITACIÓN SARAMPIÓN CLÍNICA Y VACUNACIÓN 17FEBRERO 2026.pptx'
prs = Presentation(file_path)

print("Slides:")
for i, slide in enumerate(prs.slides):
    text = ""
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            text += shape.text + " "
    if 'DEFINICIÓN OPERACIONAL' in text.upper() or 'DEFINICION OPERACIONAL' in text.upper():
        print(f"Found DEFINICION OPERACIONAL at slide {i}")
    if 'DIAGNOSTICO DIFERENCIAL' in text.upper() or 'DIAGNÓSTICO DIFERENCIAL' in text.upper():
        print(f"Found DIAGNOSTICO DIFERENCIAL at slide {i}")
    
