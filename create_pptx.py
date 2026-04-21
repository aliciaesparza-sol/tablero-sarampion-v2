import os
import random
import pptx
from pptx.util import Inches, Pt

def create_presentation():
    # Load template
    template_path = r"C:\Users\aicil\.gemini\antigravity\scratch\plantilla.pptx"
    prs = pptx.Presentation(template_path)
    
    # Clean existing slides if needed? Actually we are reopening the original template, so we don't need to clean.
    
    # Title Slide
    try:
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "CÁMARA JURISDICCIONAL NUM 3"
        subtitle.text = "Reporte de Condiciones Físicas y Estructurales"
    except Exception as e:
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

    # General Conditions Slide
    try:
        content_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(content_layout)
        title = slide.shapes.title
        if title:
            title.text = "Resumen del Estado General"
        body = slide.shapes.placeholders[1] if len(slide.shapes.placeholders) > 1 else None
        if body:
            tf = body.text_frame
            tf.text = "Durante la inspección de la Cámara Jurisdiccional Num 3 se detectaron diversas mermas críticas:"
            p = tf.add_paragraph()
            p.text = "- Presencia de humedad excesiva y filtraciones."
            p.level = 1
            p = tf.add_paragraph()
            p.text = "- Desgaste visible y mal estado en las superficies."
            p.level = 1
            p = tf.add_paragraph()
            p.text = "- Agujeros y daños estructurales severos."
            p.level = 1
    except Exception as e:
        pass

    # Find photos
    photo_dir = r"C:\Users\aicil\.gemini\antigravity\scratch\Jurisdiccion3Photos"
    files = [f for f in os.listdir(photo_dir) if f.endswith(".jpeg")]
    
    dark_photo = "WhatsApp Image 2026-04-16 at 4.25.12 AM (7).jpeg"
    general_files = [f for f in files if f != dark_photo]
    
    # To keep output consistent but varied, sort files to have a stable order
    general_files.sort()
    
    # Function to add an image slide
    def add_image_slide(title_text, desc_text, img_file1, img_file2=None):
        slide = prs.slides.add_slide(content_layout)
        if slide.shapes.title: slide.shapes.title.text = title_text
        try:
            body = slide.shapes.placeholders[1] if len(slide.shapes.placeholders) > 1 else None
            if body:
                body.text_frame.text = desc_text
                
                # Resize text area somewhat manually by placing images below it
                body.top = Inches(1.5)
                body.height = Inches(1.0)
                
            # Add images
            if img_file1:
                slide.shapes.add_picture(os.path.join(photo_dir, img_file1), Inches(1), Inches(3), width=Inches(3.5))
            if img_file2:
                slide.shapes.add_picture(os.path.join(photo_dir, img_file2), Inches(5), Inches(3), width=Inches(3.5))
        except Exception as e:
            print("Error slide:", e)

    # 1. Humedad y Filtraciones
    add_image_slide(
        "Evidencia de Humedad y Filtraciones",
        "Se observan serias manchas de humedad y escurrimientos en el interior, lo cual compromete el correcto funcionamiento de la cámara y representa un riesgo.",
        general_files[0], general_files[1]
    )
    
    # 2. Desgaste y Mal Estado
    add_image_slide(
        "Desgaste y Mal Estado General",
        "La estructura y los recubrimientos presentan un nivel de deterioro avanzado. Falta de mantenimiento preventivo evidente en diversas zonas.",
        general_files[2], general_files[3]
    )

    # 3. Oxidación y Suciedad
    add_image_slide(
        "Falta de Mantenimiento y Acumulación",
        "Se documenta la presencia material acumulado, óxido y condiciones de higiene que no cumplen con los estándares requeridos para su operación.",
        general_files[4], general_files[5]
    )
    
    # 4. Grietas y Desprendimientos
    add_image_slide(
        "Grietas y Desprendimiento de Material",
        "Paredes y techos muestran grietas profundas y desprendimiento progresivo del aislante o material de recubrimiento.",
        general_files[6], general_files[7]
    )

    # 5. Hallazgo Específico: Agujero
    slide = prs.slides.add_slide(content_layout)
    if slide.shapes.title:
        slide.shapes.title.text = "Daño Crítico: Agujero Estructural"
    try:
        body = slide.shapes.placeholders[1] if len(slide.shapes.placeholders) > 1 else None
        if body:
            tf = body.text_frame
            tf.text = "En esta zona oscura se logró captar una pequeña entrada de luz que refleja la existencia de un agujero significativo que atraviesa la estructura de la cámara de forma peligrosa."
        slide.shapes.add_picture(os.path.join(photo_dir, dark_photo), Inches(2.5), Inches(2.5), width=Inches(5))
    except Exception as e:
        print("Error black photo:", e)

    # 6. Conclusión
    try:
        slide = prs.slides.add_slide(content_layout)
        if slide.shapes.title: slide.shapes.title.text = "Conclusión y Observaciones Finales"
        body = slide.shapes.placeholders[1] if len(slide.shapes.placeholders) > 1 else None
        if body:
            tf = body.text_frame
            tf.text = "El estado integral de la Cámara Jurisdiccional Num 3 se califica como CRÍTICO."
            p = tf.add_paragraph()
            p.text = "- Urge reparación del agujero detectado para frenar pérdidas."
            p.level = 1
            p = tf.add_paragraph()
            p.text = "- Intervención inmediata para controlar el exceso de humedad."
            p.level = 1
            p = tf.add_paragraph()
            p.text = "- Se requiere un mantenimiento general para garantizar su operabilidad."
            p.level = 1
    except Exception as e:
        pass


    # Save
    out_path = r"c:\Users\aicil\OneDrive\Escritorio\Presentacion_Camara_Jurisdiccional_3.pptx"
    prs.save(out_path)
    print("Presentation saved at:", out_path)

if __name__ == "__main__":
    create_presentation()
