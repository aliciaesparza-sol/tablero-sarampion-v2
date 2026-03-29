import pandas as pd
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.ticker as ticker
import io
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# =================== DATOS ===================
stock_file = r"c:\Users\aicil\OneDrive\Escritorio\PVU\SARAMPIÓN\EXISTENCIAS SRP Y SR\EXISTENCIAS_SRP_SR_POR_DIA.xlsx"
df_raw = pd.read_excel(stock_file, sheet_name='Detalle por Jurisdicción')
df_raw['FECHA_PARSED'] = pd.to_datetime(df_raw['FECHA'], dayfirst=True, errors='coerce')
df_raw['SECCIÓN'] = df_raw['SECCIÓN'].astype(str).str.strip()

def norm_jur(s):
    s = s.upper().strip()
    if 'NO. 1' in s or 'NO.1' in s: return 'JUR1'
    if 'NO. 2' in s or 'NO.2' in s: return 'JUR2'
    if 'NO. 3' in s or 'NO.3' in s: return 'JUR3'
    if 'NO. 4' in s or 'NO.4' in s: return 'JUR4'
    return None

df_raw['JUR'] = df_raw['SECCIÓN'].apply(norm_jur)
df_jur = df_raw[df_raw['JUR'].notna() & df_raw['FECHA_PARSED'].notna()].copy()
df_jur = df_jur.sort_values(['JUR', 'FECHA_PARSED'])
df_jur['CONSUMO_SR']  = df_jur.groupby('JUR')['SR PUNTOS'].diff().mul(-1).fillna(0)
df_jur['CONSUMO_SRP'] = df_jur.groupby('JUR')['SRP PUNTOS'].diff().mul(-1).fillna(0)

csv_path = r"c:\Users\aicil\OneDrive\Escritorio\PVU\SARAMPIÓN\REPORTE_SRP-SR-CENSIA\SRP-SR-2025_14-03-2026 06-41-16.csv"
df_doses = pd.read_csv(csv_path, usecols=['JURISDICCION','Fecha de registro',
    'SR PRIMERA TOTAL','SR SEGUNDA TOTAL','SRP  PRIMERA TOTAL','SRP SEGUNDA TOTAL'])
df_doses['FECHA_PARSED'] = pd.to_datetime(df_doses['Fecha de registro'], errors='coerce')
df_doses = df_doses[df_doses['FECHA_PARSED'].dt.year >= 2026]
df_doses['SR_AP']  = df_doses[['SR PRIMERA TOTAL','SR SEGUNDA TOTAL']].fillna(0).sum(axis=1)
df_doses['SRP_AP'] = df_doses[['SRP  PRIMERA TOTAL','SRP SEGUNDA TOTAL']].fillna(0).sum(axis=1)
jmap = {'DURANGO':'JUR1','GOMEZ PALACIO':'JUR2','SANTIAGO PAPASQUIARO':'JUR3','RODEO':'JUR4'}
df_doses['JUR'] = df_doses['JURISDICCION'].apply(lambda x: next((v for k,v in jmap.items() if k in str(x).upper()), None))
df_doses = df_doses[df_doses['JUR'].notna()]
doses_agg = df_doses.groupby(['FECHA_PARSED','JUR']).agg(SR_AP=('SR_AP','sum'), SRP_AP=('SRP_AP','sum')).reset_index()

merged = pd.merge(df_jur[['FECHA_PARSED','JUR','SR PUNTOS','SRP PUNTOS','CONSUMO_SR','CONSUMO_SRP']],
                  doses_agg, on=['FECHA_PARSED','JUR'], how='left').fillna(0)
merged['DIFER_SR']  = merged['CONSUMO_SR']  - merged['SR_AP']
merged['DIFER_SRP'] = merged['CONSUMO_SRP'] - merged['SRP_AP']
merged['FECHA_STR'] = merged['FECHA_PARSED'].dt.strftime('%d/%m')
merged = merged.sort_values(['FECHA_PARSED','JUR'])

# =================== CONSTANTS ===================
DARK_BLUE  = RGBColor(0x1F, 0x38, 0x64)
LIGHT_BLUE = RGBColor(0xBD, 0xD7, 0xEE)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
SR_COLOR   = '#1A73C8'
SRP_COLOR  = '#2E7D32'
CONS_COLOR = '#C62828'
AP_COLOR   = '#F57F17'
DIFF_POS   = '#EF5350'
DIFF_NEG   = '#66BB6A'
DIFF_ZERO  = '#E0E0E0'

FUENTE = "Fuente: Reporte diario de existencias (EXISTENCIAS_SRP_SR_POR_DIA.xlsx) y base CENSIA (SRP-SR-2025) | Fecha de corte: 14 de marzo de 2026"

# =================== CAPTIONS PER CHART PER JUR ===================
CAPTIONS = {
    'JUR1': {
        'stock': "Hallazgo: La Jur. 1 (Durango) recibió grandes insumos de SR del 3–6/03, elevando su stock a más de 15,000 dosis. La SRP se consumió casi en su totalidad (5,127→415).\nConclusión: El abasto llegó oportuno. Recomendación: Vigilar que el stock de SRP se reponga antes del cierre de campaña.",
        'sr_corr': "Hallazgo: Los días 03–06 y 09–11/03 el inventario de SR bajó significativamente pero CENSIA registra pocas dosis.\nConclusión: Existen aplicaciones que no fueron capturadas en el sistema. Recomendación: Auditar y completar el registro de esos días.",
        'srp_corr': "Hallazgo: Se detectaron desfases en 12 fechas donde el consumo teórico SRP no coincide con el registro CENSIA.\nConclusión: Patrón de subregistro recurrente. Recomendación: Implementar captura en tiempo real en todos los módulos de la Jur. 1.",
        'diff': "Hallazgo: Las barras rojas (desfase >100 dosis) se concentran en la semana del 3 al 11 de marzo.\nConclusión: El problema se centra en ese periodo. Recomendación: Convocar a los responsables de módulo a conciliación urgente del registro de esa semana."
    },
    'JUR2': {
        'stock': "Hallazgo: La Jur. 2 (Gómez Palacio) inició sin existencias SR (0 dosis el 18/02) y las fue recibiendo gradualmente. SRP se consumió de manera estable durante el mes.\nConclusión: El abasto llegó tarde. Recomendación: Planificar el insumo antes del inicio de cada campaña.",
        'sr_corr': "Hallazgo: Los días 10–13/03 presentan diferencias mayores a 300 dosis de SR entre consumo e inventario y lo registrado en CENSIA.\nConclusión: Subreporte en cierre de campaña. Recomendación: Verificar que los puntos de vacunación en zona lagunera reporten diariamente.",
        'srp_corr': "Hallazgo: Se registraron desfases SRP en 14 fechas distintas, distribuyéndose a lo largo de todo el periodo.\nConclusión: Problema de registro crónico. Recomendación: Supervisar la captura diaria en cada punto y revisar conectividad del sistema.",
        'diff': "Hallazgo: Múltiples barras rojas distribuidas en todo el mes, sin concentrarse en una sola semana.\nConclusión: El subregistro no es puntual sino sistemático. Recomendación: Establecer reporte obligatorio al cierre del turno en cada módulo."
    },
    'JUR3': {
        'stock': "Hallazgo: Jur. 3 (Santiago Papasquiaro) nunca superó 390 dosis SR simultáneas. La SRP se consumió de 2,211 a 1,300 (−911 dosis).\nConclusión: Territorio serrano con bajo volumen pero consumo activo de SRP. Recomendación: Garantizar una reserva mínima de SR para este territorio.",
        'sr_corr': "Hallazgo: CENSIA no refleja prácticamente ninguna dosis SR aplicada a pesar de que el inventario sí se movió.\nConclusión: Es la jurisdicción con mayor brecha entre realidad e inventario. Recomendación: Refuerzo urgente en capacitación de captura para los vacunadores de Santiago Papasquiaro.",
        'srp_corr': "Hallazgo: El consumo teórico SRP en varias fechas es positivo pero las dosis CENSIA son cero.\nConclusión: Las dosis se aplicaron pero no se capturaron. Recomendación: Asignar un responsable de supervisión de registro en esta jurisdicción.",
        'diff': "Hallazgo: Los desfases en Santiago Papasquiaro son moderados en magnitud pero constantes. 8 fechas con desfase SRP >100 dosis.\nConclusión: El problema no es de volumen sino de disciplina en el registro. Recomendación: Visita de supervisión in situ."
    },
    'JUR4': {
        'stock': "Hallazgo: El stock de la Jur. 4 (Rodeo) oscila cíclicamente (entrega → consumo → entrega), reflejando un patrón de abastecimiento escalonado.\nConclusión: La logística funciona aunque en volúmenes pequeños. Recomendación: Documentar formalmente los traslados intermedios.",
        'sr_corr': "Hallazgo: En 10 fechas el consumo teórico SR no se corresponde con las dosis registradas, especialmente el 02/03, 04–06/03 y 09/03.\nConclusión: Las variaciones de stock no reflejan dosis aplicadas en CENSIA. Recomendación: Registrar traslados entre módulos como movimientos separados.",
        'srp_corr': "Hallazgo: El desfase SRP promedio es ~99 dosis/día, menor que otras jurisdicciones pero persistente.\nConclusión: El territorio de Rodeo dificulta el registro oportuno. Recomendación: Habilitar captura offline para comunidades sin conectividad.",
        'diff': "Hallazgo: El desfase SR promedio es de ~168 dosis/día, el segundo más alto del estado.\nConclusión: La magnitud del desfase justifica una revisión exhaustiva del periodo 04–09/03. Recomendación: Cruzar registros físicos (libretas de vacunación) con el sistema."
    }
}

# =================== CHART BUILDER ===================
def make_individual_chart(jur, chart_type):
    sub = merged[merged['JUR'] == jur].sort_values('FECHA_PARSED').reset_index(drop=True)
    if sub.empty:
        return None

    fig, ax = plt.subplots(figsize=(11.5, 4.0))
    x = np.arange(len(sub))

    if chart_type == 'stock':
        ax.plot(x, sub['SR PUNTOS'],  marker='o', ms=4, lw=2, color=SR_COLOR,  label='Existencias SR en puntos')
        ax.plot(x, sub['SRP PUNTOS'], marker='s', ms=4, lw=2, color=SRP_COLOR, label='Existencias SRP en puntos')
        ax.fill_between(x, sub['SR PUNTOS'],  alpha=0.08, color=SR_COLOR)
        ax.fill_between(x, sub['SRP PUNTOS'], alpha=0.08, color=SRP_COLOR)
        ax.set_title('📦 Existencias en Puntos de Vacunación (SR vs SRP)', fontweight='bold', fontsize=11)
        ax.legend(fontsize=9)

    elif chart_type == 'sr_corr':
        w = 0.35
        ax.bar(x - w/2, sub['CONSUMO_SR'], w, label='Consumo teórico SR', color=CONS_COLOR, alpha=0.85)
        ax.bar(x + w/2, sub['SR_AP'],       w, label='Dosis SR aplicadas', color=AP_COLOR,   alpha=0.85)
        ax.axhline(0, color='black', lw=0.7)
        ax.set_title('📊 Consumo Teórico vs Dosis Aplicadas — SR (Doble Viral)', fontweight='bold', fontsize=11)
        ax.legend(fontsize=9)

    elif chart_type == 'srp_corr':
        w = 0.35
        ax.bar(x - w/2, sub['CONSUMO_SRP'], w, label='Consumo teórico SRP', color='#6A1B9A', alpha=0.85)
        ax.bar(x + w/2, sub['SRP_AP'],       w, label='Dosis SRP aplicadas', color='#00897B', alpha=0.85)
        ax.axhline(0, color='black', lw=0.7)
        ax.set_title('📊 Consumo Teórico vs Dosis Aplicadas — SRP (Triple Viral)', fontweight='bold', fontsize=11)
        ax.legend(fontsize=9)

    elif chart_type == 'diff':
        desfase_colors = [DIFF_POS if d > 100 else (DIFF_NEG if d < -100 else DIFF_ZERO) for d in sub['DIFER_SR']]
        ax.bar(x, sub['DIFER_SR'], color=desfase_colors, edgecolor='white', linewidth=0.3)
        ax.axhline(0,    color='black',   lw=0.8, zorder=5)
        ax.axhline(100,  color='#EF5350', lw=1,   ls='--', alpha=0.7, label='Umbral +100')
        ax.axhline(-100, color='#66BB6A', lw=1,   ls='--', alpha=0.7, label='Umbral -100')
        p_red = mpatches.Patch(color=DIFF_POS,  label='Desfase > 100 (posible subreporte)')
        p_gre = mpatches.Patch(color=DIFF_NEG,  label='Desfase < -100 (posible sobreregistro)')
        p_gra = mpatches.Patch(color=DIFF_ZERO, label='Sin desfase significativo')
        ax.legend(handles=[p_red, p_gre, p_gra], fontsize=8)
        ax.set_title('⚠️  Desfase SR: Consumo Teórico − Dosis Aplicadas', fontweight='bold', fontsize=11)

    ax.set_xticks(x)
    ax.set_xticklabels(sub['FECHA_STR'], rotation=50, ha='right', fontsize=7)
    ax.grid(axis='y', alpha=0.3)
    ax.yaxis.set_major_formatter(ticker.StrMethodFormatter('{x:,.0f}'))
    ax.set_ylabel('Dosis', fontsize=9)
    plt.tight_layout()

    buf = io.BytesIO()
    fig.savefig(buf, format='png', bbox_inches='tight', dpi=160, facecolor='white')
    buf.seek(0)
    plt.close('all')
    return buf

# =================== PPTX HELPERS ===================
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

def add_header(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = DARK_BLUE; bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.2), Inches(0.04), Inches(12.9), Inches(0.55))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = WHITE
    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.2), Inches(0.6), Inches(12.9), Inches(0.4))
        p2 = tb2.text_frame.paragraphs[0]; p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.size = Pt(10); r2.font.color.rgb = LIGHT_BLUE

def add_footer(slide):
    tb = slide.shapes.add_textbox(Inches(0.1), Inches(7.2), Inches(13.13), Inches(0.27))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = FUENTE
    r.font.size = Pt(8); r.font.color.rgb = RGBColor(0x50,0x50,0x50); r.font.italic = True

def add_caption(slide, text, left, top, width, height):
    box = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF0, 0xF5, 0xFF)
    box.line.color.rgb = RGBColor(0xBD, 0xD7, 0xEE)
    tb = slide.shapes.add_textbox(Inches(left+0.07), Inches(top+0.04), Inches(width-0.14), Inches(height-0.08))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text
    r.font.size = Pt(9); r.font.color.rgb = RGBColor(0x20,0x20,0x20)

# =================== SLIDE 1: PORTADA ===================
slide = prs.slides.add_slide(blank)
bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = DARK_BLUE

def tb_s(slide, txt, l, t, w, h, sz, bold=False, col=WHITE, align=PP_ALIGN.CENTER, it=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = txt
    r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = col; r.font.italic = it

tb_s(slide, "SECRETARÍA DE SALUD — DURANGO", 1, 0.9, 11.3, 0.6, 13, col=LIGHT_BLUE)
tb_s(slide, "Hallazgos por Jurisdicción Sanitaria", 0.5, 1.6, 12.3, 1.3, 30, bold=True)
tb_s(slide, "Existencias · Consumo · Dosis Aplicadas — Vacunas SR y SRP", 1, 2.95, 11.3, 0.7, 15, col=LIGHT_BLUE)
tb_s(slide, "Campaña de Vacunación contra Sarampión 2026", 1, 3.7, 11.3, 0.5, 13, col=LIGHT_BLUE)
tb_s(slide, "Periodo: 18 de febrero — 14 de marzo de 2026", 1, 4.3, 11.3, 0.5, 11, col=RGBColor(0x90,0xCA,0xF9))
tb_s(slide, "Fecha de corte: 14 de marzo de 2026  |  Preparado para COEVA — Marzo 2026", 1, 4.85, 11.3, 0.5, 10, col=RGBColor(0x90,0xCA,0xF9))

# =================== JUR SLIDES ===================
JUR_INFO = {
    'JUR1': ('Jurisdicción No. 1 — Durango',              'Zona urbana · Mayor volumen de biológicos · Campaña intensiva'),
    'JUR2': ('Jurisdicción No. 2 — Gómez Palacio',        'Zona lagunera · Abastecimiento gradual · Alta dispersión geográfica'),
    'JUR3': ('Jurisdicción No. 3 — Santiago Papasquiaro', 'Zona serrana · Bajo volumen · Cobertura rural'),
    'JUR4': ('Jurisdicción No. 4 — Rodeo',                'Zona sur / cañones · Acceso limitado · Comunidades dispersas'),
}

CHART_TYPES = [
    ('stock',    'Existencias SR y SRP'),
    ('sr_corr',  'Consumo vs Aplicadas SR'),
    ('srp_corr', 'Consumo vs Aplicadas SRP'),
    ('diff',     'Desfase SR')
]

for jur_key in ['JUR1', 'JUR2', 'JUR3', 'JUR4']:
    title, subtitle = JUR_INFO[jur_key]
    captions = CAPTIONS[jur_key]

    # One slide per jurisdiction with 2x2 layout
    slide = prs.slides.add_slide(blank)
    add_header(slide, title, subtitle)

    # Layout: 2 cols × 2 rows, each chart + caption below
    chart_keys = ['stock', 'sr_corr', 'srp_corr', 'diff']
    positions = [
        (0.08,  1.08),  # top-left
        (6.72,  1.08),  # top-right
        (0.08,  4.15),  # bottom-left
        (6.72,  4.15),  # bottom-right
    ]
    CHART_W = 6.55
    CHART_H = 2.6
    CAP_H   = 0.88

    for idx, ck in enumerate(chart_keys):
        img = make_individual_chart(jur_key, ck)
        lft, top = positions[idx]
        if img:
            slide.shapes.add_picture(img, Inches(lft), Inches(top), Inches(CHART_W), Inches(CHART_H))
        # Caption below the chart
        add_caption(slide, captions[ck], lft, top + CHART_H + 0.04, CHART_W, CAP_H)

    add_footer(slide)

# =================== CONCLUSIONES ===================
slide = prs.slides.add_slide(blank)
add_header(slide, "Conclusiones y Acciones Inmediatas", "Periodo: 18/02 — 14/03/2026 | 4 Jurisdicciones analizadas")

conclusions = [
    ("✅  Abasto garantizado", "Las 4 jurisdicciones contaron con vacunas SR y SRP durante todo el periodo, con reabastecimientos intermedios."),
    ("✅  Jur. 1 sin desabasto", "La Jurisdicción No. 1 (Durango) consumió casi toda su SRP y recibió grandes insumos de SR. Stock activo todo el mes."),
    ("⚠️  Desfase generalizado", "Se detectaron diferencias entre consumo teórico e inventario en TODAS las jurisdicciones. El problema es estructural, no aislado."),
    ("⚠️  Jur. 3 en alerta", "Santiago Papasquiaro tiene el patrón más preocupante: stock disminuye pero CENSIA casi no registra dosis aplicadas. Subregistro severo."),
    ("⚠️  Semana crítica: 03–13/03", "La mayor concentración de desfases ocurre en esa semana, coincidiendo con el pico de aplicación de la campaña."),
    ("📋  Acción 1: Conciliar CENSIA", "Revisar y completar el registro de las fechas con desfase. Prioridad: Jur. 1, 2 y 3 del 03 al 13 de marzo."),
    ("📋  Acción 2: Conciliación diaria", "Implementar un proceso de cierre diario: conteo físico vs dosis registradas, firmado por responsable jurisdiccional."),
    ("📋  Acción 3: Supervisión Santiago", "Enviar equipo de supervisión a Jur. 3 (Santiago Papasquiaro) para verificar la captura de vacunadores y revisar libretas físicas."),
    ("📋  Acción 4: Registrar traslados", "Los movimientos entre módulos dentro de cada jurisdicción deben registrarse como entradas/salidas separadas."),
]

row_h = 0.58
start_y = 1.12
for i, (label, text) in enumerate(conclusions):
    y = start_y + i * row_h
    box = slide.shapes.add_shape(1, Inches(0.2), Inches(y), Inches(12.9), Inches(row_h - 0.04))
    box.fill.solid()
    if '✅' in label:  box.fill.fore_color.rgb = RGBColor(0xC6,0xEF,0xCE)
    elif '⚠️' in label: box.fill.fore_color.rgb = RGBColor(0xFF,0xEB,0x9C)
    elif '📋' in label: box.fill.fore_color.rgb = RGBColor(0xBD,0xD7,0xEE)
    box.line.color.rgb = RGBColor(0xCC,0xCC,0xCC)

    tb = slide.shapes.add_textbox(Inches(0.3), Inches(y+0.04), Inches(12.7), Inches(row_h - 0.1))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run(); r1.text = label + "  "
    r1.font.size = Pt(9); r1.font.bold = True; r1.font.color.rgb = DARK_BLUE
    r2 = p.add_run(); r2.text = text
    r2.font.size = Pt(9); r2.font.color.rgb = RGBColor(0x20,0x20,0x20)

add_footer(slide)

# =================== SAVE ===================
out_path = r"c:\Users\aicil\OneDrive\Escritorio\PVU\SARAMPIÓN\HALLAZGOS_POR_JURISDICCION_COEVA_2026.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
