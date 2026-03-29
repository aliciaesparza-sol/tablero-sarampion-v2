from pptx import Presentation

file_path = r'c:\Users\aicil\OneDrive\Escritorio\PVU\COEVA\PRESENTACIONES\CAPACITACIÓN SARAMPIÓN CLÍNICA Y VACUNACIÓN 17FEBRERO 2026.pptx'
prs = Presentation(file_path)

with open('all_contents_final.txt', 'w', encoding='utf-8') as f:
    for i, slide in enumerate(prs.slides):
        f.write(f"--- S{i} ---\n")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                f.write(f"[{shape.name}] {shape.text.replace('\n', ' ')}\n")
        f.write("\n")
