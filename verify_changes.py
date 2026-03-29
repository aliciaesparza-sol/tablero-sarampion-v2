from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

file_path = r'c:\Users\aicil\OneDrive\Escritorio\PVU\COEVA\PRESENTACIONES\CAPACITACIÓN SARAMPIÓN CLÍNICA Y VACUNACIÓN 17FEBRERO 2026.pptx'
prs = Presentation(file_path)

print("--- VERIFICATION ---")

# Check Slide 9
slide_9 = prs.slides[9]
found_source = False
for shape in slide_9.shapes:
    if hasattr(shape, "text") and "Fuente: OPS" in shape.text:
        found_source = True
        print(f"Slide 9: Found 'Fuente: OPS' in shape '{shape.name}'")
if not found_source:
    print("Slide 9: 'Fuente: OPS' NOT FOUND")

# Check Slide 18
slide_18 = prs.slides[18]
found_table = False
for shape in slide_18.shapes:
    if shape.has_table:
        found_table = True
        table = shape.table
        print(f"Slide 18: Found table with {len(table.rows)} rows and {len(table.columns)} columns")
        # Check first cell of header and data
        print(f"  Header[0,0]: '{table.cell(0,0).text}'")
        print(f"  Data[1,0]: '{table.cell(1,0).text}'")
if not found_table:
    print("Slide 18: Table NOT FOUND")
