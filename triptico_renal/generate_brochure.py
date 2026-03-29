import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_pptx():
    prs = Presentation()
    
    # Set slide size to A4 Landscape (approx)
    prs.slide_width = Inches(11.69)
    prs.slide_height = Inches(8.27)

    def add_column_text(slide, x, y, width, title, bullets, title_size=20, bullet_size=12, color=(0, 64, 128)):
        txBox = slide.shapes.add_textbox(x, y, width, Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(title_size)
        p.font.color.rgb = RGBColor(*color)
        p.alignment = PP_ALIGN.CENTER

        curr_y = y + Inches(0.6)
        for bullet in bullets:
            txBox2 = slide.shapes.add_textbox(x + Inches(0.2), curr_y, width - Inches(0.4), Inches(0.3))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = f"• {bullet}"
            p2.font.size = Pt(bullet_size)
            p2.font.color.rgb = RGBColor(50, 50, 50)
            curr_y += Inches(0.4)

    # --- SLIDE 1: INTERIOR ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6]) # blank layout
    
    col_width = prs.slide_width / 3
    
    # Column 1: ¿QUÉ ES? + CAUSAS
    add_column_text(slide1, Inches(0.1), Inches(0.5), col_width - Inches(0.2), "¿QUÉ ES?", [
        "Es la pérdida progresiva e irreversible de la función renal.",
        "Impide eliminar desechos y líquidos adecuadamente."
    ])
    add_column_text(slide1, Inches(0.1), Inches(2.5), col_width - Inches(0.2), "PRINCIPALES CAUSAS", [
        "Diabetes", "Hipertensión", "Infecciones urinarias", "Medicamentos nefrotóxicos"
    ])
    slide1.shapes.add_picture("lifestyle.png", Inches(0.5), Inches(5.5), width=Inches(3))

    # Column 2: SIGNOS Y SÍNTOMAS
    add_column_text(slide1, col_width + Inches(0.1), Inches(0.5), col_width - Inches(0.2), "SIGNOS Y SÍNTOMAS", [
        "Cansancio y debilidad", "Cambios en la orina", "Edema en pies o rostro", "Náuseas", "Comezón", "Palidez"
    ])
    
    # Column 3: CUIDADOS DE ENFERMERÍA
    add_column_text(slide1, 2*col_width + Inches(0.1), Inches(0.5), col_width - Inches(0.2), "CUIDADOS DE ENFERMERÍA", [
        "Control de signos vitales", "Balance hídrico estricto", "Peso diario", 
        "Administración correcta de medicamentos", "Vigilancia de complicaciones", "Prevención de infecciones"
    ])
    slide1.shapes.add_picture("nursing.png", 2*col_width + Inches(0.5), Inches(5.5), width=Inches(3))

    # --- SLIDE 2: EXTERIOR ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Column 1: EDUCACIÓN AL PACIENTE (Solapa interna)
    add_column_text(slide2, Inches(0.1), Inches(0.5), col_width - Inches(0.2), "EDUCACIÓN AL PACIENTE", [
        "Apego al tratamiento", "Dieta baja en sal", "No automedicarse", 
        "Asistir a consultas y laboratorios", "Reconocer signos de alarma"
    ])

    # Column 2: ¿CUÁNDO ACUDIR? (Contraportada)
    add_column_text(slide2, col_width + Inches(0.1), Inches(0.5), col_width - Inches(0.2), "¿CUÁNDO ACUDIR AL\nSERVICIO DE SALUD?", [
        "Disminución importante de la orina", "Dificultad para respirar", "Hinchazón excesiva", 
        "Dolor intenso", "Náuseas o vómitos persistentes"
    ], title_size=18)
    
    footer = slide2.shapes.add_textbox(col_width + Inches(0.1), Inches(7), col_width - Inches(0.2), Inches(0.5))
    footer.text_frame.text = "La atención oportuna mejora la calidad de vida."
    footer.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Column 3: PORTADA
    # Background for cover
    slide2.shapes.add_picture("cover.png", 2*col_width, 0, width=col_width, height=prs.slide_height)
    
    # Cover Text Box
    cover_box = slide2.shapes.add_textbox(2*col_width + Inches(0.2), Inches(0.5), col_width - Inches(0.4), Inches(3.5))
    tf = cover_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "ENFERMEDAD RENAL\nCRÓNICA"
    p.font.bold = True
    p.font.size = Pt(30)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "Cuidados de Enfermería"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(200, 200, 200)
    p2.alignment = PP_ALIGN.CENTER
    
    p3 = tf.add_paragraph()
    p3.text = "\nPrevención • Atención • Educación"
    p3.font.size = Pt(12)
    p3.font.color.rgb = RGBColor(255, 255, 255)
    p3.alignment = PP_ALIGN.CENTER

    p_name = tf.add_paragraph()
    p_name.text = "\nAlumna:\nMyriam Machado Arreola"
    p_name.font.size = Pt(14)
    p_name.font.color.rgb = RGBColor(255, 255, 255)
    p_name.alignment = PP_ALIGN.CENTER
    
    p_serv = tf.add_paragraph()
    p_serv.text = "\nServicio de Enfermería"
    p_serv.font.size = Pt(12)
    p_serv.font.color.rgb = RGBColor(200, 200, 200)
    p_serv.alignment = PP_ALIGN.CENTER

    prs.save("Triptico_Renal.pptx")

if __name__ == "__main__":
    create_pptx()
