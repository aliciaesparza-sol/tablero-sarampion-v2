import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, PatternFill
import os

# Configuración de estética
sns.set_theme(style="whitegrid")
plt.rcParams['figure.figsize'] = (12, 6)
COLOR_2025 = '#2E86C1'
COLOR_2026 = '#E67E22'
COLOR_HIGHLIGHT = '#C0392B'

CSV_PATH = r"c:\Descargas_SRP\SRP-SR-2025_20-04-2026 03-35-04.csv"
TEMPLATE_PATH = r"c:\Users\aicil\OneDrive\Escritorio\plantilla presentacion ppt.pptx"
OUTPUT_DIR = r"C:\Users\aicil\.gemini\antigravity\scratch"
EXCEL_OUTPUT = os.path.join(OUTPUT_DIR, "reporte_vacunacion_SRP_SR_v4.xlsx")
PPT_OUTPUT = os.path.join(OUTPUT_DIR, "presentacion_vacunacion_SRP_SR_v4.pptx")
CHART_DAILY = os.path.join(OUTPUT_DIR, "chart_daily_doses.png")
CHART_CUMULATIVE = os.path.join(OUTPUT_DIR, "chart_cumulative_doses.png")
CHART_MONTHLY = os.path.join(OUTPUT_DIR, "chart_monthly_doses.png")
CHART_TOP2026 = os.path.join(OUTPUT_DIR, "chart_top_days_2026.png")

MONTH_NAMES = {1: 'Ene', 2: 'Feb', 3: 'Mar', 4: 'Abr', 5: 'May', 6: 'Jun', 
               7: 'Jul', 8: 'Ago', 9: 'Sep', 10: 'Oct', 11: 'Nov', 12: 'Dic'}

def process_data():
    print("Cargando y procesando datos...")
    df = pd.read_csv(CSV_PATH)
    dose_cols = ['SRP  PRIMERA TOTAL', 'SRP SEGUNDA TOTAL', 'SR PRIMERA TOTAL', 'SR SEGUNDA TOTAL']
    df['Dosis Totales'] = df[dose_cols].sum(axis=1)
    df['Fecha_Dt'] = pd.to_datetime(df['Fecha de registro'])
    df['Fecha'] = df['Fecha_Dt'].dt.date
    df['Año'] = df['Fecha_Dt'].dt.year
    df['Mes'] = df['Fecha_Dt'].dt.month
    df['Dia_Anio'] = df['Fecha_Dt'].dt.dayofyear
    return df.groupby(['Fecha', 'Año', 'Mes', 'Dia_Anio'])['Dosis Totales'].sum().reset_index()

def create_charts(daily_doses):
    print("Generando gráficos para v4...")
    # Diario
    plt.figure(figsize=(14, 7)); sns.lineplot(data=daily_doses, x='Dia_Anio', y='Dosis Totales', hue='Año', palette={2025: COLOR_2025, 2026: COLOR_2026})
    plt.title('Comparativa Diaria de Vacunación (2025 vs 2026)', fontsize=16, fontweight='bold'); plt.tight_layout(); plt.savefig(CHART_DAILY, dpi=300); plt.close()
    
    # Mensual
    plt.figure(figsize=(14, 7)); m_sum = daily_doses.groupby(['Año', 'Mes'])['Dosis Totales'].sum().reset_index(); m_sum['Nombre_Mes'] = m_sum['Mes'].map(MONTH_NAMES)
    sns.barplot(data=m_sum, x='Nombre_Mes', y='Dosis Totales', hue='Año', palette={2025: COLOR_2025, 2026: COLOR_2026})
    plt.title('Resumen Mensual de Dosis Aplicadas', fontsize=16, fontweight='bold'); plt.tight_layout(); plt.savefig(CHART_MONTHLY, dpi=300); plt.close()

    # Top
    top_2026 = daily_doses[daily_doses['Año'] == 2026].sort_values('Dosis Totales', ascending=False).head(10).copy()
    top_2026['Fecha_Str'] = top_2026['Fecha'].astype(str)
    plt.figure(figsize=(14, 7)); sns.barplot(data=top_2026, x='Dosis Totales', y='Fecha_Str', color=COLOR_HIGHLIGHT)
    plt.title('Top 10 Días de Mayor Registro en 2026', fontsize=16, fontweight='bold'); plt.tight_layout(); plt.savefig(CHART_TOP2026, dpi=300); plt.close()

    # Acumulado
    plt.figure(figsize=(14, 7))
    for yr, col in [(2025, COLOR_2025), (2026, COLOR_2026)]:
        dy = daily_doses[daily_doses['Año'] == yr].copy().sort_values('Dia_Anio')
        dy['Acumulado'] = dy['Dosis Totales'].cumsum()
        sns.lineplot(data=dy, x='Dia_Anio', y='Acumulado', label=f'Acumulado {yr}', color=col, linewidth=3)
    plt.title('Avance Acumulado de Vacunación', fontsize=16, fontweight='bold'); plt.tight_layout(); plt.savefig(CHART_CUMULATIVE, dpi=300); plt.close()

def create_excel(daily_doses):
    print("Generando Excel v4...")
    with pd.ExcelWriter(EXCEL_OUTPUT, engine='openpyxl') as writer:
        daily_doses.to_excel(writer, sheet_name='Dosis_Diarias', index=False)
        top = daily_doses[daily_doses['Año'] == 2026].sort_values('Dosis Totales', ascending=False).head(20)
        top.to_excel(writer, sheet_name='Top_2026', index=False)

def add_template_slide(prs, title, img, text, layout_idx=5):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    # Intentar asignar título si existe el placeholder
    if slide.shapes.title:
        slide.shapes.title.text = title
    else:
        # Si no hay placeholder de título, agregar un textbox arriba
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    # Imagen
    slide.shapes.add_picture(img, Inches(0.5), Inches(1.5), width=Inches(9))
    
    # Texto explicativo (Pie de foto)
    tb_pie = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1))
    tf = tb_pie.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.CENTER

def create_ppt(daily_doses):
    print("Generando PowerPoint v4 con plantilla institucional...")
    try:
        prs = Presentation(TEMPLATE_PATH)
    except:
        print("No se encontró la plantilla, usando una nueva.")
        prs = Presentation()
    
    # 1. Diapositiva de Portada (Layout 0)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if slide.shapes.title:
        slide.shapes.title.text = "Análisis de Vacunación SRP-SR"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Informe Ejecutivo de Metas y Productividad\nComparativa 2025 - 2026\nFuente: CENSIA / SRP-SR"

    # Diapositivas con gráficas (Layout 5: Solo título)
    add_template_slide(prs, "Mejores Jornadas de Vacunación (Top 10 - 2026)", CHART_TOP2026, 
                       "Análisis de picos operativos durante el año actual.")
    
    add_template_slide(prs, "Comparativa Mensual de Productividad", CHART_MONTHLY, 
                       "Contraste del volumen registrado mes a mes entre periodos anuales.")
    
    add_template_slide(prs, "Avance Histórico Diario (2025 vs 2026)", CHART_DAILY, 
                       "Monitoreo de la consistencia diaria en la aplicación y registro de dosis.")
    
    add_template_slide(prs, "Cumplimiento de Metas: Desempeño Acumulado", CHART_CUMULATIVE, 
                       "Visibilidad de la trayectoria anual vs el año previo para evaluación de metas.")
    
    # Tabla final
    slide_table = prs.slides.add_slide(prs.slide_layouts[1])
    if slide_table.shapes.title:
        slide_table.shapes.title.text = "Resumen Consolidado de Dosis"
    totals = daily_doses.groupby('Año')['Dosis Totales'].sum().reset_index()
    rows, cols = 3, 2
    table = slide_table.shapes.add_table(rows, cols, Inches(2), Inches(2.5), Inches(6), Inches(1.5)).table
    table.cell(0, 0).text = "Año Gestión"; table.cell(0, 1).text = "Dosis Totales"
    for i, row in totals.iterrows():
        table.cell(i+1, 0).text = str(int(row['Año'])); table.cell(i+1, 1).text = "{:,.0f}".format(row['Dosis Totales'])

    prs.save(PPT_OUTPUT)

if __name__ == "__main__":
    try:
        data = process_data()
        create_charts(data)
        create_excel(data)
        create_ppt(data)
        print(f"Éxito: Reporte v4 generado en {OUTPUT_DIR}")
    except Exception as e:
        print(f"Error: {e}")
