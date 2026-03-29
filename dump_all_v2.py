from pptx import Presentation

file_path = r'c:\Users\aicil\OneDrive\Escritorio\PVU\COEVA\PRESENTACIONES\CAPACITACIÓN SARAMPIÓN CLÍNICA Y VACUNACIÓN 17FEBRERO 2026.pptx'
prs = Presentation(file_path)

for i, slide in enumerate(prs.slides):
    title = slide.shapes.title.text if slide.shapes.title else "No Title"
    content = " | ".join([shape.text.replace("\n", " ")[:50] for shape in slide.shapes if hasattr(shape, "text") and shape != slide.shapes.title])
    print(f"Slide {i}: Title: {title} | Content: {content[:150]}")
