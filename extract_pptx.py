from pptx import Presentation
import json

def extract_pptx_info(file_path):
    prs = Presentation(file_path)
    slides_info = []
    
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        
        slides_info.append({
            "slide_index": i + 1,
            "text": slide_text
        })
    
    return slides_info

if __name__ == "__main__":
    info = extract_pptx_info("Vacunacion_SRP_SR_Presentacion_.pptx")
    with open("presentation_content.json", "w", encoding="utf-8") as f:
        json.dump(info, f, ensure_ascii=False, indent=2)
    print("Content extracted to presentation_content.json")
