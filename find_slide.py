import pptx
import sys

ppt_path = r"c:\Users\aicil\OneDrive\Escritorio\PVU\COEVA\PRESENTACIONES\CAPACITACIÓN SARAMPIÓN CLÍNICA Y VACUNACIÓN 17FEBRERO 2026.pptx"

try:
    prs = pptx.Presentation(ppt_path)
    for i, slide in enumerate(prs.slides):
        print(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)
                if "manejo" in shape.text.lower() and "sospechoso" in shape.text.lower():
                    print(f"\n\n*** FOUND MATCH ON SLIDE {i+1} ***\n\n")
except Exception as e:
    print(f"Error: {e}")
